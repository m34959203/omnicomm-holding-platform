"""Сборка клиентского отчёта в PowerPoint (.pptx), Excel (.xlsx), HTML и PDF.

ТЗ §8 — структура и стиль презентации (8 слайдов, светлый корпоративный
дизайн, 16:9). ТЗ §9 — текстовые формулировки без обвинений. ТЗ §10 —
выгрузка очищенных/рассчитанных данных в Excel и опциональный PDF.

Бизнес-инварианты (CLAUDE.md, ТЗ §7):
  * колонка «возможные сливы топлива» НИКОГДА не попадает ни в .pptx, ни в .xlsx;
  * аномалии подаются только как «требуют проверки», без обвинительных слов;
  * перерасход не утверждаем без согласованных норм;
  * читаемые шрифты, ничего не вылезает за пределы слайда.

Модуль работает только с единой моделью (`FleetReport`), от источника
данных (API/Excel) не зависит. Графики приходят готовыми PNG-путями от
`charts.build_charts` под ключами: 'mileage', 'fuel_per_100km',
'fuel_idle', 'speeding'.
"""

from __future__ import annotations

import base64
import logging
import math
import os
import shutil
import subprocess
from html import escape
from typing import Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from . import analytics, config, svg_charts, vehicle_types
from .models import FleetReport, Severity, VehicleMetrics

logger = logging.getLogger(__name__)


# --- Палитра и типографика (ТЗ §8) -------------------------------------------
# Согласовано визуально с charts.py: спокойный корпоративный набор.
COLOR_PRIMARY = RGBColor(0x2F, 0x5C, 0x8F)   # синий — заголовки, заливка таблицы
COLOR_TEXT = RGBColor(0x2B, 0x2B, 0x2B)      # тёмно-серый — основной текст
COLOR_ACCENT = RGBColor(0xC8, 0x89, 0x3F)    # тёплый акцент — KPI, выделения
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)        # белый фон слайда
COLOR_BG_SOFT = RGBColor(0xF4, 0xF6, 0xF9)   # очень светло-серый — подложки
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_MUTED = RGBColor(0x8A, 0x8A, 0x8A)     # приглушённый — футер, сноски

FONT_FAMILY = "Calibri"
SIZE_TITLE = Pt(34)        # заголовки слайдов (ТЗ: 28–36pt)
SIZE_TITLE_HERO = Pt(36)   # титул
SIZE_SUBTITLE = Pt(18)
SIZE_BODY = Pt(16)         # тело (ТЗ: 14–18pt)
SIZE_SMALL = Pt(14)
SIZE_FOOTER = Pt(10)
SIZE_KPI_VALUE = Pt(30)
SIZE_KPI_LABEL = Pt(13)

# Размер слайда 16:9 в EMU (1 дюйм = 914400 EMU → 13.333" × 7.5").
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(610000)       # ~0.5" поля, чтобы контент не вылезал за слайд

# Ограничение строк на таблицу: подобрано так, чтобы таблица гарантированно
# влезала между шапкой и примечанием/футером при кегле тела 10pt в одну
# строку (раньше 15 строк с переносом длинных имён вылезали за футер).
MAX_ROWS_PER_TABLE = 10
SIZE_TABLE_BODY = Pt(10)
# Макс. длина текста в первой колонке таблиц (имена ТС/площадки/причины) —
# чтобы строка не переносилась на 2-ю и не раздувала высоту строки.
MAX_CELL_LABEL = 22

# Текстовые инварианты — формулировки без обвинений (ТЗ §7, §9).
TXT_NO_DATA = "нет данных"
TXT_REVIEW = "требует проверки"
TXT_NORMS_DISCLAIMER = (
    "Перерасход нельзя считать без утверждённых норм расхода: значения ниже — "
    "относительные, для приоритизации проверки, а не вывод о нарушении."
)
TXT_CHART_UNAVAILABLE = "График недоступен"


# --- Низкоуровневые helper'ы текста ------------------------------------------

def _style_run(run, *, size: Pt, color: RGBColor, bold: bool = False) -> None:
    """Единообразно оформить run: шрифт, размер, цвет, начертание."""
    run.font.name = FONT_FAMILY
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    size: Pt = SIZE_BODY,
    color: RGBColor = COLOR_TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Добавить текстовый блок. Включаем word_wrap, чтобы текст не вылезал."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _style_run(run, size=size, color=color, bold=bold)
    return box


def _blank_slide(prs: Presentation):
    """Пустой слайд (layout 6) — строим всё вручную для контроля верстки."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Светлый фон на каждый слайд (официальный корпоративный стиль).
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    return slide


def _add_footer(slide, report: FleetReport) -> None:
    """Футер: клиент + период мелким текстом — фирменная подпись каждого слайда."""
    text = f"{report.client_name}   ·   {report.period.human()}"
    _add_textbox(
        slide,
        MARGIN,
        Emu(SLIDE_H - 430000),
        Emu(SLIDE_W - 2 * MARGIN),
        Emu(330000),
        text,
        size=SIZE_FOOTER,
        color=COLOR_MUTED,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_slide_header(slide, title: str) -> None:
    """Крупный заголовок слайда + тонкая акцентная линия под ним.

    Размер заголовка подбирается по длине: длинные заголовки уменьшаются,
    чтобы остаться в одну строку и не наезжать на акцентную линию ниже.
    """
    n = len(title)
    size = SIZE_TITLE if n <= 30 else (Pt(28) if n <= 46 else Pt(24))
    _add_textbox(
        slide,
        MARGIN,
        Emu(360000),
        Emu(SLIDE_W - 2 * MARGIN),
        Emu(900000),
        title,
        size=size,
        color=COLOR_PRIMARY,
        bold=True,
    )
    # Акцентная линия-разделитель под заголовком.
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        MARGIN,
        Emu(1230000),
        Emu(1500000),
        Emu(45000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


# --- Форматирование значений (единый стиль чисел в отчёте) -------------------

def _fmt(value: Optional[float], unit: str = "", digits: int = 0) -> str:
    """Число с разделителем тысяч и единицей; None → прочерк."""
    if value is None:
        return "—"
    if digits == 0:
        body = f"{value:,.0f}".replace(",", " ")
    else:
        body = f"{value:,.{digits}f}".replace(",", " ")
    return f"{body} {unit}".strip()


def _clip(text: str, maxlen: int = MAX_CELL_LABEL) -> str:
    """Обрезать длинный текст ячейки с многоточием, чтобы он не переносился
    на 2-ю строку и не раздувал высоту строки таблицы (→ переполнение)."""
    text = (text or "").strip()
    return text if len(text) <= maxlen else text[:maxlen - 1].rstrip() + "…"


def _col_capacity(col_w_emu: int, font_pt: float) -> int:
    """Сколько символов помещается в колонку в ОДНУ строку при кегле font_pt."""
    col_px = col_w_emu / 9525 - 8
    char_px = font_pt * 96 / 72 * 0.52
    return max(6, int(col_px / max(1.0, char_px)))


# --- Шаблоны слайдов (helper'ы из ТЗ) ----------------------------------------

def _add_title_slide(prs: Presentation, report: FleetReport) -> None:
    """Слайд 1 — титул: тема, клиент, период, число ТС, 3–4 ключевых KPI."""
    slide = _blank_slide(prs)
    kpi = report.kpi

    _add_textbox(
        slide, MARGIN, Emu(900000), Emu(SLIDE_W - 2 * MARGIN), Emu(1300000),
        "Аналитический отчёт по автопарку",
        size=SIZE_TITLE_HERO, color=COLOR_PRIMARY, bold=True,
    )
    _add_textbox(
        slide, MARGIN, Emu(2150000), Emu(SLIDE_W - 2 * MARGIN), Emu(560000),
        f"{report.client_name}",
        size=SIZE_SUBTITLE, color=COLOR_TEXT, bold=True,
    )
    _add_textbox(
        slide, MARGIN, Emu(2620000), Emu(SLIDE_W - 2 * MARGIN), Emu(480000),
        f"Период: {report.period.human()}   ·   ТС в анализе: {kpi.vehicles_total} "
        f"(с данными: {kpi.vehicles_with_data})",
        size=SIZE_BODY, color=COLOR_MUTED,
    )

    # 3–4 ключевых KPI карточками в ряд.
    cards = [
        ("Общий пробег", _fmt(kpi.total_mileage_km, "км")),
        ("Общий расход", _fmt(kpi.total_fuel_l, "л")),
        ("Средний расход", _fmt(kpi.weighted_fuel_per_100km, "л/100км", 1)),
        ("Моточасы", _fmt(kpi.total_engine_hours, "ч", 1)),
    ]
    n = len(cards)
    gap = Emu(220000)
    total_w = SLIDE_W - 2 * MARGIN
    card_w = Emu((int(total_w) - int(gap) * (n - 1)) // n)
    card_h = Emu(1380000)
    top = Emu(3300000)
    for i, (label, value) in enumerate(cards):
        left = Emu(int(MARGIN) + i * (int(card_w) + int(gap)))
        card = slide.shapes.add_shape(1, left, top, card_w, card_h)  # RECTANGLE
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_SOFT
        card.line.color.rgb = COLOR_PRIMARY
        card.line.width = Pt(0.75)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(120000)
        tf.margin_right = Emu(120000)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        _style_run(r1, size=SIZE_KPI_VALUE, color=COLOR_ACCENT, bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        _style_run(r2, size=SIZE_KPI_LABEL, color=COLOR_TEXT)

    # Плашка-вывод внизу — закрывает «мёртвую» нижнюю половину титула и
    # сразу подаёт главный денежный смысл отчёта.
    _add_title_takeaway(slide, report)

    _add_footer(slide, report)


def _add_title_takeaway(slide, report: FleetReport) -> None:
    """Итоговая плашка под KPI на титуле: главный инсайт периода."""
    kpi = report.kpi
    story = _money_story(report)
    if story:
        headline = "Главный резерв периода"
        body = (f"{_fmt(story['idle_cost'])} ₸ на холостом ходу — "
                f"{story['idle_share'] * 100:.0f}% всех денег на топливо; "
                f"в пересчёте на год ≈ {_fmt(story['idle_annual'])} ₸.")
    elif kpi.total_fuel_cost > 0:
        headline = "Стоимость топлива за период"
        body = (f"{_fmt(kpi.total_fuel_cost)} ₸, из них на простоях с "
                f"работающим двигателем {_fmt(kpi.idle_fuel_cost)} ₸.")
    else:
        headline = "Что в отчёте"
        body = (f"{kpi.vehicles_with_data} ТС с данными; холостой ход "
                f"{kpi.idle_hours_share * 100:.0f}% моточасов. Детализация, "
                f"деньги и приоритеты — на следующих слайдах.")

    box_top = 5000000
    box_h = 1180000
    plate = slide.shapes.add_shape(1, MARGIN, Emu(box_top),
                                   Emu(SLIDE_W - 2 * MARGIN), Emu(box_h))
    plate.fill.solid()
    plate.fill.fore_color.rgb = COLOR_BG_SOFT
    plate.line.fill.background()
    tf = plate.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(260000)
    tf.margin_right = Emu(260000)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = headline
    _style_run(r1, size=SIZE_SMALL, color=COLOR_MUTED, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = body
    _style_run(r2, size=SIZE_SUBTITLE, color=COLOR_PRIMARY, bold=True)


def _paginate_bullets(bullets: list[str], avail_first: int, avail_rest: int,
                      font_pt: float, space_pt: float) -> list[list[str]]:
    """Разбить пункты по слайдам так, чтобы каждая страница влезала в рабочую
    зону (с учётом переноса длинных пунктов) — вместо вылета за футер/обрезки."""
    width_px = (int(SLIDE_W) - 2 * int(MARGIN)) / EMU_PX
    fpx = int(font_pt * 96 / 72)
    line_h = int(fpx * 1.3) * EMU_PX
    cpl = max(12, int((width_px - 12) / (font_pt * 96 / 72 * 0.52)))
    space = int(space_pt * 12700)
    pages: list[list[str]] = []
    cur: list[str] = []
    used, avail = 0, avail_first
    for b in bullets:
        lines = max(1, math.ceil((len(b) + 2) / cpl))
        h = lines * line_h + space
        if cur and used + h > avail:
            pages.append(cur)
            cur, used, avail = [], 0, avail_rest
        cur.append(b)
        used += h
    if cur:
        pages.append(cur)
    return pages or [[]]


def _render_bullets_page(prs, report, title, page, *, intro, size, space):
    """Один слайд со списком тезисов (страница пагинации)."""
    slide = _blank_slide(prs)
    _add_slide_header(slide, title)
    top = 1450000
    if intro:
        _add_textbox(slide, MARGIN, Emu(top), Emu(SLIDE_W - 2 * MARGIN),
                     Emu(560000), intro, size=SIZE_BODY, color=COLOR_MUTED)
        top += 620000
    box = slide.shapes.add_textbox(
        MARGIN, Emu(top), Emu(SLIDE_W - 2 * MARGIN), Emu(BODY_BOTTOM - top))
    tf = box.text_frame
    tf.word_wrap = True
    # Страховка: если оценка пагинации ошиблась, PowerPoint ужмёт текст под рамку.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(page):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = space
        marker = p.add_run()
        marker.text = "— "
        _style_run(marker, size=size, color=COLOR_ACCENT, bold=True)
        run = p.add_run()
        run.text = text
        _style_run(run, size=size, color=COLOR_TEXT)
    _add_footer(slide, report)


def _add_bullets_slide(
    prs: Presentation,
    report: FleetReport,
    title: str,
    bullets: list[str],
    *,
    intro: Optional[str] = None,
) -> None:
    """Слайд(ы) со списком тезисов. При большом числе пунктов — пагинация
    на несколько слайдов с суффиксом (N/M), как у таблиц; контент не теряется
    и не уходит за футер."""
    safe = bullets or ["Данных для выводов недостаточно — требуется проверка источника."]
    # Кегль выбираем по объёму: меньше пунктов — крупнее (для выводов),
    # больше — компактнее (меньше страниц). Переполнение решает пагинация.
    size = SIZE_BODY if len(safe) <= 7 else SIZE_SMALL
    space = Pt(8)
    intro_h = 620000 if intro else 0
    avail_first = BODY_BOTTOM - (1450000 + intro_h)
    avail_rest = BODY_BOTTOM - 1450000
    pages = _paginate_bullets(safe, avail_first, avail_rest, size.pt, space.pt)
    m = len(pages)
    for i, page in enumerate(pages, start=1):
        ptitle = title if m == 1 else f"{title} ({i}/{m})"
        _render_bullets_page(prs, report, ptitle, page,
                             intro=intro if i == 1 else None,
                             size=size, space=space)


def _add_image_slide(
    prs: Presentation,
    report: FleetReport,
    title: str,
    image_path: Optional[str],
    note: str,
    *,
    extra_lines: Optional[list[str]] = None,
) -> None:
    """Слайд с графиком слева/по центру + текстовый вывод снизу.

    Картинка вставляется только если файл реально существует; иначе —
    плашка «График недоступен» (ТЗ: не падать на отсутствии графика).
    """
    slide = _blank_slide(prs)
    _add_slide_header(slide, title)

    img_top = 1450000
    img_left = MARGIN
    img_w = Emu(SLIDE_W - 2 * MARGIN)

    # Подпись считаем ПЕРВОЙ: её высота определяет, сколько места под график.
    # Длинные перечни не вылезают на футер — шрифт ужимается, число строк
    # ограничивается, а блок пиннится над футером.
    lines = [note] + list(extra_lines or [])
    if len(lines) > 9:  # защита от чрезмерных перечней
        lines = lines[:8] + ["…"]
    n_lines = len(lines)
    cap_size = Pt(14) if n_lines <= 4 else (Pt(12) if n_lines <= 6 else Pt(11))
    # Реальная высота подписи — с учётом переноса длинных строк (иначе
    # подпись из длинных пунктов наезжала на футер).
    cap_w_emu = int(SLIDE_W) - 2 * int(MARGIN)
    visual_lines = sum(_est_cell_lines(t, cap_w_emu, cap_size.pt, True)
                       for t in lines)
    cap_line_emu = int(cap_size.pt * 12700 * 1.35) + 30000
    cap_h = visual_lines * cap_line_emu + 120000
    cap_top = FOOTER_TOP - 90000 - cap_h
    cap_top = max(cap_top, 2900000)  # не залезать на шапку при малом числе строк

    img_h = Emu(max(1500000, cap_top - img_top - 160000))

    if image_path and os.path.exists(image_path):
        # Вписываем по высоте, центрируем по горизонтали — не выходим за слайд.
        pic = slide.shapes.add_picture(image_path, img_left, Emu(img_top),
                                       height=img_h)
        if pic.width > int(img_w):
            pic.height = None  # пересчёт по ширине
            pic.width = img_w
            pic.height = int(img_w * pic.image.size[1] / pic.image.size[0])
        pic.left = Emu(int(MARGIN) + (int(img_w) - pic.width) // 2)
        # если ужали по ширине — поднимем чуть, чтобы держать по центру зоны
        pic.top = Emu(img_top + max(0, (int(img_h) - pic.height) // 2))
    else:
        plate = slide.shapes.add_shape(1, img_left, Emu(img_top), img_w, img_h)
        plate.fill.solid()
        plate.fill.fore_color.rgb = COLOR_BG_SOFT
        plate.line.color.rgb = COLOR_MUTED
        plate.line.width = Pt(0.75)
        tf = plate.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = TXT_CHART_UNAVAILABLE
        _style_run(r, size=SIZE_SUBTITLE, color=COLOR_MUTED, bold=True)

    # Текстовый вывод под графиком.
    box = slide.shapes.add_textbox(
        MARGIN, Emu(cap_top), Emu(SLIDE_W - 2 * MARGIN), Emu(cap_h)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = text
        _style_run(run, size=cap_size, color=COLOR_TEXT, bold=(i == 0))

    _add_footer(slide, report)


# Нижняя граница рабочей зоны тела (над футером) — общая для таблиц/подписей.
BODY_BOTTOM = int(SLIDE_H) - 560000
FOOTER_TOP = int(SLIDE_H) - 430000


EMU_PX = 9525                 # 1 px при 96 dpi
MIN_DATA_ROW = 360000         # минимум высоты строки тела (воздух + ≥ контента)
MIN_HEADER_ROW = 440000


def _est_cell_lines(text: str, col_w_emu: int, font_pt: float, wrap) -> int:
    """Сколько строк займёт текст в ячейке (как при переносе PowerPoint).
    Эвристика по средней ширине символа — согласована с render-превью."""
    if not text or wrap is False:
        return 1
    col_px = col_w_emu / EMU_PX - 8
    char_px = font_pt * 96 / 72 * 0.52
    cpl = max(1, int(col_px / max(1.0, char_px)))
    return max(1, math.ceil(len(text) / cpl))


def _table_real_heights(table, body_pt: float) -> list[int]:
    """РЕАЛЬНЫЕ высоты строк (EMU) с учётом переноса — то, что форсит
    PowerPoint. Раньше высота задавалась «вслепую», и длинные имена ТС
    раздували строки за футер."""
    col_w = [c.width for c in table.columns]
    heights = []
    for ri, row in enumerate(table.rows):
        is_h = ri == 0
        fpt = SIZE_SMALL.pt if is_h else body_pt
        line_h = int(int(fpt * 96 / 72) * 1.3) * EMU_PX
        need = 0
        for ci, cell in enumerate(row.cells):
            cw = col_w[ci] if ci < len(col_w) else 40 * EMU_PX
            n = _est_cell_lines(cell.text.strip(), cw,
                                fpt, cell.text_frame.word_wrap)
            need = max(need, n * line_h + 2 * 10 * EMU_PX)
        heights.append(max(need, MIN_HEADER_ROW if is_h else MIN_DATA_ROW))
    return heights


def _fit_table(graphic, table, *, top0: int, footnote: bool,
               body_pt: float) -> tuple[int, int]:
    """Спозиционировать таблицу так, чтобы она ГАРАНТИРОВАННО влезла между
    шапкой и примечанием/футером. Вызывать ПОСЛЕ заполнения ячеек.

    Считает реальные высоты строк (wrap-aware), центрирует короткие таблицы;
    если контент не помещается — пропорционально ужимает строки под зону.
    Возвращает (top, реальная_высота) — чтобы поставить примечание под таблицей.
    """
    reserve = 520000 if footnote else 0
    avail = BODY_BOTTOM - top0 - reserve
    heights = _table_real_heights(table, body_pt)
    total = sum(heights)
    if total > avail:  # страховка (после клипа имён и лимита строк почти не нужна)
        k = avail / total
        heights = [max(int(h * k), 200000) for h in heights]
    elif total < avail:  # растягиваем строки, чтобы заполнить зону (без гигантов)
        k = avail / total
        caps = [560000 if i == 0 else 760000 for i in range(len(heights))]
        heights = [min(int(h * k), caps[i]) for i, h in enumerate(heights)]
    total = sum(heights)
    top = top0 + max(0, (avail - total) // 2)  # центрируем остаток
    graphic.top = Emu(top)
    graphic.height = Emu(total)
    for i, h in enumerate(heights):
        table.rows[i].height = Emu(h)
    return top, total


def _add_table_slide(
    prs: Presentation,
    report: FleetReport,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    footnote: Optional[str] = None,
) -> None:
    """Слайд с таблицей. Шапка — заливка фирменным синим, строки — зебра."""
    slide = _blank_slide(prs)
    _add_slide_header(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_left = MARGIN
    tbl_top = Emu(1450000)
    tbl_w = Emu(SLIDE_W - 2 * MARGIN)
    tbl_h = Emu(min(4400000, 360000 * n_rows))

    graphic = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h)
    table = graphic.table

    # Шапка.
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        _style_run(r, size=SIZE_SMALL, color=COLOR_WHITE, bold=True)

    # Тело: зебра + первый столбец слева, числа по центру. Первую колонку
    # клипуем ПО ШИРИНЕ колонки — длинные имена ТС/площадок не переносятся на
    # 2-ю строку (иначе строка раздувается и таблица вылезает за футер). Клип
    # по ширине, а не фикс-длине: узкая колонка профиля ≠ широкая колонка
    # экономики — последнюю не режем зря.
    col0_cap = _col_capacity((SLIDE_W - 2 * MARGIN) // n_cols, SIZE_TABLE_BODY.pt)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_SOFT if ri % 2 == 0 else COLOR_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = _clip(val, col0_cap) if c == 0 else val
            _style_run(r, size=SIZE_TABLE_BODY, color=COLOR_TEXT)

    # Позиционируем ПОСЛЕ заполнения — по реальным высотам строк.
    top, tbl_h_fit = _fit_table(graphic, table, top0=1450000,
                                footnote=bool(footnote),
                                body_pt=SIZE_TABLE_BODY.pt)

    if footnote:
        foot_top = min(top + tbl_h_fit + 260000, FOOTER_TOP - 360000)
        _add_textbox(
            slide, MARGIN, Emu(foot_top), Emu(SLIDE_W - 2 * MARGIN), Emu(330000),
            footnote, size=SIZE_FOOTER, color=COLOR_MUTED,
        )

    _add_footer(slide, report)


# --- Сборка контента слайдов 3–7 ---------------------------------------------

def _vehicle_table_rows(vehicles: list[VehicleMetrics]) -> list[list[str]]:
    """Строки таблицы профиля парка. БЕЗ колонки сливов (инвариант ТЗ §7)."""
    rows: list[list[str]] = []
    for v in vehicles:
        if not v.has_data:
            reason = v.no_data_reason or TXT_NO_DATA
            rows.append([v.name, TXT_NO_DATA, "", "", "", "", _clip(f"({reason})", 16)])
            continue
        # Метрика расхода выбирается ШАБЛОНОМ типа ТС (vehicle_types):
        # для техники работы на месте — л/моточас, для перевозок — л/100 км.
        prof = vehicle_types.profile(v.vehicle_type)
        if prof.primary_metric == "l_per_mh":
            rate = f"{_fmt(v.fuel_per_motorhour, '', 1)} л/мч"
        else:
            rate = _fmt(v.fuel_per_100km_calc, "", 1)
        rows.append([
            v.name,
            _fmt(v.mileage_km, "", 0),
            _fmt(v.fuel_l, "", 1),
            rate,
            _fmt(v.engine_hours, "", 1),
            _fmt(v.engine_idle_hours, "", 1),
            _fmt(v.max_speed_kmh, "", 0),
        ])
    return rows


def _add_fleet_profile_slides(prs: Presentation, report: FleetReport) -> None:
    """Слайд(ы) 3 — профиль парка. При >15 ТС дробим на несколько слайдов."""
    headers = ["ТС", "Пробег, км", "Расход, л", "Расход уд.", "Моточасы",
               "Простой, ч", "Макс. V"]
    rows = _vehicle_table_rows(report.vehicles)
    title = "Профиль автопарка"

    if not rows:
        _add_table_slide(prs, report, title, headers,
                         [["—", TXT_NO_DATA, "", "", "", "", ""]])
        return

    # Дробим на слайды, балансируя размер чанков — чтобы не было вырожденного
    # последнего слайда в 1 строку (пустой слайд + флаг пустоты).
    n_parts = math.ceil(len(rows) / MAX_ROWS_PER_TABLE)
    per = math.ceil(len(rows) / n_parts)
    chunks = [rows[i:i + per] for i in range(0, len(rows), per)]
    total_parts = len(chunks)
    for idx, chunk in enumerate(chunks, start=1):
        part_title = title if total_parts == 1 else f"{title} ({idx}/{total_parts})"
        footnote = (
            f"ТС с пометкой «{TXT_NO_DATA}» исключены из агрегатов. "
            "Колонка расхода без движения вынесена на отдельный слайд."
        )
        _add_table_slide(prs, report, part_title, headers, chunk, footnote=footnote)


def _high_fuel_vehicles(report: FleetReport, top_n: int = 5) -> list[str]:
    """ТС с наибольшим относительным расходом — для приоритизации проверки.

    Стационарную спецтехнику (погрузчики/экскаваторы) НЕ ранжируем по л/100 км:
    при околонулевом пробеге показатель вырастает до бессмысленных значений
    (напр. погрузчик «1087 л/100км»). Она измеряется в л/моточас. (Cowork-ревью)
    """
    with_rate = [
        (v.name, v.fuel_per_100km_calc)
        for v in report.vehicles
        if v.has_data and v.fuel_per_100km_calc is not None and not v.is_stationary
    ]
    with_rate.sort(key=lambda x: x[1], reverse=True)
    return [f"{name}: {rate:.1f} л/100км" for name, rate in with_rate[:top_n]]


def _is_stationary_equipment(v: VehicleMetrics) -> bool:
    """Неподвижная спецтехника (метрика л/моточас): по паспорту типа или кинематике."""
    if vehicle_types.profile(v.vehicle_type).primary_metric == "l_per_mh":
        return True
    return v.is_stationary


def _has_stationary_equipment(report: FleetReport) -> bool:
    return any(_is_stationary_equipment(v) for v in report.vehicles if v.has_data)


def _spec_fuel_per_mh_lines(report: FleetReport, top_n: int = 5) -> list[str]:
    """Топ спецтехники по расходу л/моточас — для подписи слайда."""
    rows = [
        (v.name, v.fuel_per_motorhour)
        for v in report.vehicles
        if v.has_data and _is_stationary_equipment(v)
        and v.fuel_per_motorhour is not None and v.fuel_per_motorhour > 0
    ]
    rows.sort(key=lambda x: x[1], reverse=True)
    return [f"{name}: {rate:.1f} л/моточас" for name, rate in rows[:top_n]]


_TREND_LABELS = {
    "total_mileage_km": "пробег",
    "total_fuel_l": "расход топлива",
    "weighted_fuel_per_100km": "средний расход",
    "total_engine_hours": "моточасы",
    "idle_hours_share": "доля холостого хода",
    "total_fuel_cost": "стоимость топлива",
}


def _trend_lines(report: FleetReport) -> list[str]:
    """Человекочитаемые строки динамики период-к-периоду со стрелками."""
    if not report.trends:
        return []
    parts: list[str] = []
    for key, label in _TREND_LABELS.items():
        delta = report.trends.get(key)
        if delta is None or abs(delta) < 0.1:
            continue
        arrow = "▲" if delta > 0 else "▼"
        parts.append(f"{label} {arrow} {abs(delta):.0f}%")
    if not parts:
        return []
    return [f"Динамика к прошлому периоду: {'; '.join(parts)}."]


def _review_anomalies(report: FleetReport, limit: int = 8) -> list[str]:
    """Список аномалий severity=REVIEW — строго «требуют проверки» (ТЗ §7)."""
    items: list[str] = []
    for v in report.vehicles:
        for a in v.anomalies:
            if a.severity == Severity.REVIEW:
                # message уже содержит «— требует проверки» (validator), не дублируем.
                msg = a.message
                if TXT_REVIEW not in msg.lower():
                    msg = f"{msg} — {TXT_REVIEW}"
                items.append(f"{v.name}: {msg}")
    return items[:limit]


# --- Модуль «Работа на погрузке» ---------------------------------------------

def _loading_relevant(report: FleetReport) -> bool:
    """Есть ли смысл показывать раздел погрузки (хоть один сигнал/точка)."""
    return any(
        (v.loading_method in ("sensor", "sensor_zero", "rpm", "gps", "geozone"))
        or v.loading_points
        for v in report.vehicles if v.has_data
    )


def _fleet_fuel_per_mh_stationary(report: FleetReport) -> Optional[float]:
    """Расход на моточас работы стоя по парку, л/мч (дуальная норма, идея Антона)."""
    fuel = sum(v.idle_fuel_wo_move_l or 0 for v in report.vehicles if v.has_data)
    hours = sum(v.work_no_move_hours or 0 for v in report.vehicles if v.has_data)
    return round(fuel / hours, 1) if hours else None


_METHOD_BADGE = {
    "sensor": "датчик надстройки",
    "sensor_zero": "датчик: не грузил",
    "rpm": "по оборотам ≈",
    "gps": "по GPS-маршруту ≈",
    "geozone": "по геозонам-площадкам ≈",
    "none": "нет сигнала",
}


def _loading_rows(report: FleetReport) -> list[list[str]]:
    """Строки таблицы экономики погрузки по ТС (только с сигналом)."""
    rows = []
    for v in report.vehicles:
        if not v.has_data or v.loading_method in (None, "none"):
            continue
        util = v.loading_utilization
        rows.append([
            v.name,
            _fmt(v.loading_hours, "", 1) if v.loading_hours is not None else "—",
            f"{util * 100:.0f}%" if util is not None else "—",
            _fmt(v.loading_fuel_l, "", 1) if v.loading_fuel_l is not None else "—",
            _fmt(v.loading_fuel_per_mh, "", 1) if v.loading_fuel_per_mh is not None else "—",
            _METHOD_BADGE.get(v.loading_method, v.loading_method or "—"),
        ])
    return rows


def _has_geozone_visits(report: FleetReport) -> bool:
    return any(p.source == "geozone" for v in report.vehicles for p in v.loading_points)


def _geozone_journal_rows(report: FleetReport, top_n: int = 15) -> list[list[str]]:
    """Журнал по площадкам (геозонам): площадка → визитов, время, топливо, ₸."""
    from collections import defaultdict
    agg = defaultdict(lambda: {"n": 0, "dur": 0.0, "fuel": 0.0})
    price = report.kpi.fuel_price_kzt
    for v in report.vehicles:
        for p in v.loading_points:
            if p.source != "geozone":
                continue
            a = agg[p.name or "—"]
            a["n"] += 1
            a["dur"] += p.duration_s
            a["fuel"] += p.fuel_l or 0.0
    rows = []
    for name, a in sorted(agg.items(), key=lambda x: -x[1]["n"])[:top_n]:
        cost = a["fuel"] * price if price else 0
        rows.append([
            str(name), str(a["n"]), _fmt(a["dur"] / 3600, "ч", 1),
            _fmt(a["fuel"], "л", 1), _fmt(cost) if price else "—",
        ])
    return rows


def _add_loading_slides(prs: Presentation, report: FleetReport,
                        charts: dict[str, str]) -> None:
    """Слайды раздела «Работа на погрузке»: раскладка, дуальная норма, таблица, карта."""
    kpi = report.kpi
    mh_stationary = _fleet_fuel_per_mh_stationary(report)
    fact_h = kpi.total_loading_hours_sensor
    est_h = kpi.total_loading_hours_estimate

    # Слайд: раскладка погрузка/простой + ключевые цифры.
    lines = [
        f"Время погрузки: {_fmt(fact_h, 'ч', 1)} по датчику (факт)"
        + (f" + {_fmt(est_h, 'ч', 1)} по оценке (обороты/GPS-маршрут, ≈)" if est_h else "") + ".",
        f"Доля полезной работы из времени стоя: {kpi.fleet_loading_utilization * 100:.0f}%.",
    ]
    if mh_stationary is not None:
        lines.append(
            f"Расход на моточас работы стоя (подъём контейнеров): "
            f"{_fmt(mh_stationary, 'л/моточас', 1)} — топливо списывается по "
            f"моточасам, а не только по пробегу."
        )
    if kpi.total_unproductive_fuel_cost > 0:
        lines.append(
            f"Потери на непродуктивном простое: {_fmt(kpi.total_unproductive_fuel_cost)} ₸ "
            f"({_fmt(kpi.total_unproductive_fuel_l, 'л', 1)})."
        )
    lines.append(
        f"Датчик надстройки настроен у {kpi.vehicles_with_loading_sensor} из "
        f"{kpi.vehicles_with_data} ТС; у остальных — оценка по оборотам/GPS-маршруту (≈)."
    )
    _add_image_slide(
        prs, report, "Работа на погрузке",
        charts.get("loading_split"),
        "Разделение работы двигателя на стоянке: полезная погрузка vs простой.",
        extra_lines=lines,
    )

    # Слайд: таблица экономики погрузки по ТС.
    rows = _loading_rows(report)
    if rows:
        headers = ["ТС", "Погрузка, ч", "Полезно %", "Топл., л", "л/моточас", "Источник"]
        _add_table_slide(
            prs, report, "Экономика погрузки по машинам", headers, rows[:MAX_ROWS_PER_TABLE],
            footnote="«по оборотам ≈» — оценка без датчика надстройки; "
                     "топливо погрузки точно только при датчике.",
        )

    # Журнал обслуживания площадок (только при заведённых геозонах). GPS-карта
    # точек погрузки убрана по решению заказчика.
    if _has_geozone_visits(report):
        rows = _geozone_journal_rows(report)
        headers = ["Площадка", "Визитов", "Время", "Топливо", "₸"]
        total = sum(len(v.loading_points) for v in report.vehicles)
        _add_table_slide(
            prs, report, "Журнал обслуживания площадок", headers, rows,
            footnote=f"Всего визитов в геозоны-площадки: {total}. Источник — "
                     "геозоны Omnicomm (точные площадки клиента).")


def _fleet_composition(report: FleetReport) -> list[tuple[str, int, str]]:
    """Состав парка по типам: (название типа, кол-во, специфика метрики)."""
    from collections import Counter
    c = Counter(v.vehicle_type or "other" for v in report.vehicles if v.has_data)
    out = []
    for key, n in c.most_common():
        prof = vehicle_types.profile(key)
        out.append((prof.label, n, prof.note))
    return out


def _add_classification_slide(prs: Presentation, report: FleetReport) -> None:
    """Слайд «Классификация техники»: состав парка по типам + специфика анализа."""
    comp = _fleet_composition(report)
    bullets = [f"{label} — {n} ед.: {note}" for label, n, note in comp]
    _add_bullets_slide(
        prs, report, "Классификация техники",
        bullets or ["Типы техники не определены."],
        intro="Парк разнесён по типам — для каждого свой основной параметр "
              "расхода (л/100 км для перевозок, л/моточас для работы на месте).",
    )


def _html_classification_section(report: FleetReport) -> str:
    comp = _fleet_composition(report)
    if not comp:
        return ""
    items = "".join(
        f"<li><b>{escape(label)}</b> — {n} ед.: {escape(note)}</li>"
        for label, n, note in comp
    )
    return f"""
    <section>
      <h2>Классификация техники</h2>
      <p class="muted">Для каждого типа — свой основной параметр расхода
         (л/100 км для перевозок, л/моточас для работы на месте).</p>
      <ul>{items}</ul>
    </section>
    """


def _financial_lines(report: FleetReport) -> list[str]:
    """Обогащённые строки финансового раздела (итог, год, удельно, чувствительность,
    парето, бейдж достоверности). Общая логика для pptx и html."""
    kpi = report.kpi
    price = kpi.fuel_price_kzt
    idle_pct = kpi.fuel_idle_share * 100
    annual_cost = analytics.annualize(kpi.total_fuel_cost, report.period)
    annual_savings = analytics.annualize(kpi.potential_savings, report.period)
    badge = "оценка (нет датчиков надстройки)" if kpi.savings_is_estimate else "по датчику"
    lines = [
        f"Денежный итог: топливо {_fmt(kpi.total_fuel_cost)} ₸ за период → на простоях "
        f"{_fmt(kpi.idle_fuel_cost)} ₸ ({idle_pct:.0f}%) → достижимая экономия до "
        f"{_fmt(kpi.potential_savings)} ₸ [{badge}].",
        f"Годовая проекция: бюджет топлива ~{_fmt(annual_cost)} ₸/год; "
        f"экономия на простоях до ~{_fmt(annual_savings)} ₸/год.",
        f"Удельно: {_fmt(kpi.fuel_cost_per_km, '₸/км', 1)}, "
        f"{_fmt(kpi.fuel_cost_per_mh, '₸/моточас', 1)}.",
    ]
    # Чувствительность к цене ГСМ (±5%): сумма линейна по ₸/л.
    if price > 0:
        lo, hi = kpi.total_fuel_cost * 0.95, kpi.total_fuel_cost * 1.05
        lines.append(
            f"Чувствительность к цене ГСМ (±5% от {_fmt(price)} ₸/л): бюджет "
            f"{_fmt(lo)}–{_fmt(hi)} ₸.")
    # Парето: где деньги (топ-3 ТС по ₸-потерям).
    top = analytics.rank_money_loss(report.vehicles, price, top_n=3)
    if top:
        parts = [f"{t['name']} ({_fmt(t['loss'])} ₸, {t['share'] * 100:.0f}%)" for t in top]
        lines.append("Где деньги (топ потерь): " + "; ".join(parts) + ".")
    return lines


_OVERRUN_BASIS = {"combined": "пробег+моточасы", "100km": "л/100км", "mh": "л/моточас"}


def _norms_rows(report: FleetReport) -> list[list[str]]:
    """Строки таблицы перерасхода/экономии по ТС (только с нормой)."""
    rows = []
    over = [v for v in report.vehicles if v.has_data and v.overrun_l is not None]
    over.sort(key=lambda v: v.overrun_l, reverse=True)  # сначала перерасход
    for v in over:
        sign = "перерасход" if v.overrun_l > 0 else "экономия"
        cost = v.overrun_cost_kzt
        rows.append([
            v.name,
            _fmt(v.norm_l_per_100km, "", 1) if v.norm_l_per_100km else "—",
            _fmt(v.norm_l_per_mh, "", 1) if v.norm_l_per_mh else "—",
            f"{'+' if v.overrun_l > 0 else ''}{_fmt(v.overrun_l, 'л', 1)}",
            (f"{'+' if (cost or 0) > 0 else ''}{_fmt(cost)} ₸") if cost is not None else "—",
            sign,
        ])
    return rows


def _html_norms_section(report: FleetReport, charts: Optional[dict] = None) -> str:
    """HTML-раздел перерасхода/экономии по нормам (или пусто, если норм нет)."""
    kpi = report.kpi
    if kpi.vehicles_with_norm <= 0:
        return ""
    rows = _norms_rows(report)
    trow = "".join(
        "<tr>" + "".join(f"<td>{escape(str(c))}</td>" for c in r) + "</tr>"
        for r in rows
    )
    net = kpi.total_overrun_cost - kpi.total_economy_cost
    rating = svg_charts.norm_rating(report)
    fig = (f'<figure role="img" aria-label="Перерасход и экономия к норме">{rating}</figure>'
           if rating else _html_figure(charts or {}, "overrun", "Отклонение от нормы по ТС"))
    return f"""
    <section>
      <h2>Перерасход / экономия по нормам</h2>
      {fig}
      <p>Нормы заданы по {kpi.vehicles_with_norm} ТС, перерасход у
         <b>{kpi.vehicles_over_norm}</b>. Перерасход:
         <b>{_fmt(kpi.total_overrun_l, 'л', 1)}</b> ({_fmt(kpi.total_overrun_cost)} ₸);
         экономия: {_fmt(kpi.total_economy_l, 'л', 1)} ({_fmt(kpi.total_economy_cost)} ₸).
         Сальдо по парку: <b>{'+' if net > 0 else ''}{_fmt(net)} ₸</b>.</p>
      <table class="profile"><tr><th>ТС</th><th>Норма л/100км</th><th>Норма л/мч</th>
        <th>Отклонение</th><th>₸</th><th>Итог</th></tr>{trow}</table>
      <p class="muted">Перерасход рассчитан только по ТС с заданными нормами; для
         мусоровозов учитывается и пробег (л/100км), и работа стоя (л/моточас).{
         " Применён зимний коэффициент к нормам (+10%)." if report.season == "winter" else ""}</p>
    </section>
    """


def _add_norms_slide(prs: Presentation, report: FleetReport) -> None:
    """Слайд «Перерасход / экономия по нормам» — сигнал руководителю."""
    kpi = report.kpi
    headers = ["ТС", "Норма л/100км", "Норма л/мч", "Отклонение", "₸", "Итог"]
    rows = _norms_rows(report)
    net_cost = kpi.total_overrun_cost - kpi.total_economy_cost
    season_note = (" Применён зимний коэффициент к нормам (+10%)."
                   if report.season == "winter" else "")
    summary = (
        f"Нормы заданы по {kpi.vehicles_with_norm} ТС; перерасход у "
        f"{kpi.vehicles_over_norm}. Суммарный перерасход "
        f"{_fmt(kpi.total_overrun_l, 'л', 1)} ({_fmt(kpi.total_overrun_cost)} ₸), "
        f"экономия {_fmt(kpi.total_economy_l, 'л', 1)} ({_fmt(kpi.total_economy_cost)} ₸). "
        f"Сальдо: {'+' if net_cost > 0 else ''}{_fmt(net_cost)} ₸.{season_note}"
    )
    _add_table_slide(
        prs, report, "Перерасход / экономия по нормам", headers,
        rows[:MAX_ROWS_PER_TABLE], footnote=summary,
    )


def _html_signals_section(report: FleetReport) -> str:
    """HTML: сигналы, топ проблемных ТС, «что если», бенчмарк."""
    parts = []
    if report.alerts:
        parts.append("<p><b>Сигналы:</b></p>" + _html_list(report.alerts))
    top = [c for c in report.scorecard if c["score"] > 0][:6]
    if top:
        rows = "".join(
            f"<tr><td>{escape(c['name'])}</td><td>{c['score']:.0f}</td>"
            f"<td>{escape(c['reasons'])}</td></tr>" for c in top)
        parts.append("<p><b>Топ внимания (балл проблемности):</b></p>"
                     "<table class='profile'><tr><th>ТС</th><th>Балл</th><th>Причины</th></tr>"
                     + rows + "</table>")
    whatif = [s for s in report.whatif if s["saved_kzt"]]
    if whatif:
        items = "".join(
            f"<li>Сократить простой на {s['cut'] * 100:.0f}% → экономия "
            f"<b>{_fmt(s['saved_kzt'])} ₸</b> ({_fmt(s['saved_l'], 'л', 1)})</li>"
            for s in whatif)
        parts.append("<p><b>Потенциал экономии (что если):</b></p><ul>" + items + "</ul>")
    b = report.benchmark
    if b.get("metrics"):
        items = "".join(
            f"<li>{escape(m['label'])}: у вас <b>{m['mine']}</b>, в среднем {m['peers_avg']} — "
            f"{'лучше' if m['better'] else 'хуже'}</li>" for m in b["metrics"].values())
        parts.append(f"<p><b>Сравнение со средним по {b['peers']} паркам:</b></p><ul>"
                     + items + "</ul>")
    if not parts:
        return ""
    return f"<section><h2>Сигналы, приоритеты и потенциал</h2>{''.join(parts)}</section>"


def _money_story(report: FleetReport) -> Optional[dict]:
    """Синтез нарратива для руководителя: боль → куда уходят деньги (с методом
    расчёта) → что делать. Возвращает структуру для HTML- и pptx-рендера.

    Принцип: один острый показатель боли (холостой ход — факт, не оценка),
    ранжированные «утечки» с человекочитаемым методом в каждой строке (без
    суммирования пересекающихся срезов), и конкретные действия с деньгами
    и именами ТС. None — если нет цены ГСМ (без денег истории нет).
    """
    from omnicomm_report import economics as econ_mod

    kpi = report.kpi
    if not kpi.total_fuel_cost or kpi.total_fuel_cost <= 0:
        return None
    price = kpi.fuel_price_kzt
    idle_cost = kpi.idle_fuel_cost
    idle_hours = kpi.total_idle_hours
    idle_share = kpi.fuel_idle_share  # доля топлива без движения, 0..1
    moving_cost = max(0.0, kpi.total_fuel_cost - idle_cost)
    idle_annual = analytics.annualize(idle_cost, report.period)

    eco = econ_mod.build_economics(report)
    wear = next((b for b in eco.buckets if b.key == "wear"), None)

    # Утечки — три РАЗНЫХ среза, каждый со своим методом; не суммируем.
    leaks: list[dict] = []
    if idle_cost > 0:
        leaks.append({
            "label": "Холостой ход",
            "kzt": idle_cost,
            "method": (f"{idle_share * 100:.0f}% топлива сожжено на стоянке — "
                       f"двигатель работал {idle_hours:,.0f} ч без движения; "
                       f"литры без движения × {price:.0f} ₸/л").replace(",", " "),
            "tag": "факт",
        })
    if kpi.vehicles_over_norm > 0 and kpi.total_overrun_cost > 0:
        leaks.append({
            "label": "Перерасход к нормам",
            "kzt": kpi.total_overrun_cost,
            "method": (f"{kpi.vehicles_over_norm} ТС жгут сверх нормы; "
                       "(факт − норма) × цена ГСМ"),
            "tag": "нормы",
        })
    if wear and wear.existing_kzt > 0:
        leaks.append({
            "label": "Скрытый износ от простоя",
            "kzt": wear.existing_kzt,
            "method": ("простой ускоряет ТО: часы простоя × "
                       f"{config.IDLE_WEAR_KM_PER_HOUR:.0f} км износа × ставка ТО"),
            "tag": "≈ оценка",
        })
    leaks.sort(key=lambda x: x["kzt"], reverse=True)

    # Действия — конкретные, с деньгами и именами ТС.
    actions: list[dict] = []
    whatif30 = next((s for s in reversed(report.whatif) if s.get("saved_kzt")), None)
    idle_worst = [n for n, _ in eco.worst_vehicles[:3]]
    if idle_worst:
        ret = (f" Возврат до {_fmt(whatif30['saved_kzt'])} ₸ за период "
               "при сокращении простоя на 30%." if whatif30 else "")
        actions.append({
            "title": "Запустить программу против холостого хода",
            "detail": f"Начать с лидеров простоя: {', '.join(idle_worst)}.{ret}",
        })
    over_v = sorted((v for v in report.vehicles
                     if v.has_data and (v.overrun_cost_kzt or 0) > 0),
                    key=lambda v: v.overrun_cost_kzt, reverse=True)[:3]
    if over_v:
        names = "; ".join(f"{v.name} (+{_fmt(v.overrun_cost_kzt)} ₸)" for v in over_v)
        actions.append({
            "title": f"Разобрать перерасход у {kpi.vehicles_over_norm} машин",
            "detail": f"Лидеры: {names}.",
        })
    if kpi.vehicles_with_norm > 0:
        actions.append({
            "title": "Утвердить нормы расхода по паспортам ТС",
            "detail": ("Сейчас могут стоять единые ориентировочные нормы — "
                       "из-за этого перерасход индикативный, не контрактный. "
                       "Реальные паспортные нормы делают сумму перерасхода "
                       "юридически значимой для списания ГСМ."),
        })
    dark = kpi.vehicles_total - kpi.vehicles_with_data
    mute = sum(1 for v in report.vehicles
               if any(a.code == "zero_fuel_with_activity" for a in v.anomalies))
    if dark or mute:
        bits = []
        if dark:
            bits.append(f"{dark} ТС без данных за период")
        if mute:
            bits.append(f"{mute} ТС с молчащим датчиком топлива")
        actions.append({
            "title": "Восстановить полноту телеметрии",
            "detail": "Проверить: " + "; ".join(bits)
                      + " — без них картина по парку неполная.",
        })

    return {
        "spend_total": kpi.total_fuel_cost,
        "idle_cost": idle_cost,
        "idle_share": idle_share,
        "idle_hours": idle_hours,
        "idle_annual": idle_annual,
        "moving_cost": moving_cost,
        "leaks": leaks,
        "actions": actions,
    }


def _html_money_story(report: FleetReport) -> str:
    """HTML: одностраничник-нарратив «Главное» в начале отчёта."""
    s = _money_story(report)
    if not s:
        return ""
    pct = s["idle_share"] * 100
    leak_rows = "".join(
        f'<div class="leak">'
        f'<div class="leak-rank">{i}</div>'
        f'<div class="leak-body"><div class="leak-top">'
        f'<span class="leak-label">{escape(lk["label"])}'
        f'<span class="leak-tag tag-{("est" if "≈" in lk["tag"] else "fact" if lk["tag"]=="факт" else "norm")}">'
        f'{escape(lk["tag"])}</span></span>'
        f'<span class="leak-kzt">{_fmt(lk["kzt"])} ₸</span></div>'
        f'<div class="leak-method">{escape(lk["method"])}</div></div></div>'
        for i, lk in enumerate(s["leaks"], 1))
    action_items = "".join(
        f'<li><b>{escape(a["title"])}.</b> {escape(a["detail"])}</li>'
        for a in s["actions"])
    return f"""
    <section class="story">
      <div class="story-eyebrow">Главное за период</div>
      <div class="pain">
        <div class="pain-num">{_fmt(s["idle_cost"])} ₸</div>
        <div class="pain-txt">сожжено на холостом ходу — это
          <b>{pct:.0f}% всех денег на топливо</b>. Двигатель работал стоя
          {_fmt(s["idle_hours"], 'ч')}. В пересчёте на год —
          <b>≈ {_fmt(s["idle_annual"])} ₸</b> уходит в стоянку.</div>
      </div>
      <div class="story-h">Куда уходят деньги</div>
      <div class="leaks">{leak_rows}</div>
      <p class="story-note">Цена ГСМ — средняя по календарю за период.
        «факт» — из телеметрии; «нормы» — относительно заданных норм;
        «≈ оценка» — по отраслевому коэффициенту. Срезы разные и не
        складываются в один итог.</p>
      <div class="story-h">Что делать</div>
      <ol class="actions">{action_items}</ol>
    </section>"""


def _html_economics_section(report: FleetReport) -> str:
    """HTML: корзины денег + COI — зеркало слайда «Экономический эффект»."""
    from omnicomm_report import economics as econ_mod

    eco = econ_mod.build_economics(report)
    if not eco.buckets or eco.total_existing_kzt <= 0:
        return ""
    rows = "".join(
        (lambda pre:
         f"<tr><td>{escape(b.label)}</td>"
         f"<td>{(pre + _fmt(b.existing_kzt)) if b.existing_kzt > 0 else '—'}</td>"
         f"<td>{(pre + _fmt(b.potential_kzt)) if b.potential_kzt > 0 else '—'}</td>"
         f"<td class='muted'>{escape(b.note)}</td></tr>")("≈ " if b.is_estimate else "")
        for b in eco.buckets)
    total = (f"<tr><td><b>ИТОГО</b></td><td><b>{_fmt(eco.total_existing_kzt)}</b></td>"
             f"<td><b>{_fmt(eco.total_potential_kzt)}</b></td><td></td></tr>")
    coi = ""
    if eco.period_days >= 7:
        coi = (f"<p>Оценка потерь без программы (экстраполяция периода): "
               f"<b>≈ {_fmt(eco.coi_monthly_kzt)} ₸/мес</b> "
               f"(≈ {_fmt(eco.coi_annual_kzt)} ₸/год).</p>")
    worst = ""
    if eco.worst_vehicles:
        items = "".join(f"<li>{escape(n)}: потенциал <b>{_fmt(k)} ₸</b> за период</li>"
                        for n, k in eco.worst_vehicles)
        worst = (f"<p><b>Первоочередные ТС</b> (холостой ход выше медианы парка "
                 f"{eco.median_idle_share * 100:.0f}%):</p><ul>{items}</ul>"
                 "<p class='disclaimer'>Значения — для приоритизации проверки, "
                 "не вывод о нарушениях.</p>")
    return f"""
    <section>
      <h2>Экономический эффект: оценка потерь и потенциала</h2>
      <table class='profile'>
        <tr><th>Корзина</th><th>Потери за период, ₸</th>
            <th>Потенциал возврата, ₸</th><th>Основание</th></tr>
        {rows}{total}
      </table>
      {coi}{worst}
      <p class="muted">«≈» — оценка по отраслевым коэффициентам (не измерение);
         ₸/мес — экстраполяция периода. Измеренный результат программы —
         в разделе «Счётчик экономии».</p>
    </section>"""


def _html_savings_section(report: FleetReport, charts: dict[str, str]) -> str:
    """HTML: счётчик подтверждённой экономии — зеркало слайда (Ф2)."""
    s = report.savings or {}
    entry = s.get("period")
    if not entry:
        return ""
    cum_kzt, cum_l, n = s.get("cumulative_kzt", 0), s.get("cumulative_l", 0), s.get("entries_count", 0)
    b, rates = s.get("baseline", {}), s.get("baseline", {}).get("rates", {})
    comp = entry.get("components", {})
    comp_bits = "; ".join(
        f"{label} {comp[key].get('saved_l', 0):+.0f} л"
        for key, label in (("idle", "холостой ход"), ("moving", "движение"),
                           ("stationary", "спецтехника")) if key in comp)
    word = "экономия" if cum_kzt >= 0 else "перерасход к эталону"
    sign = "экономия" if entry["saved_l"] >= 0 else "перерасход к эталону"
    cls = "good" if cum_kzt >= 0 else "warn"
    src_periods = int(b.get("source_periods", 0) or 0)
    thin = ("<p class='disclaimer'>Baseline заморожен из одного периода — "
            "оценка предварительная; точность растёт с историей.</p>"
            if src_periods < 2 else "")
    return f"""
    <section>
      <h2>Счётчик экономии против baseline</h2>
      <div class="kpis"><div class="card card-{cls}">
        <div class="card-val">{_fmt(abs(cum_kzt))} ₸</div>
        <div class="card-lbl">Накоплено за {n} пер. — {word}
          ({_fmt(abs(cum_l), 'л')}) против baseline</div>
      </div></div>
      {thin}
      {_html_figure(charts, "savings", "Накопленная экономия против baseline")}
      <p>Период {escape(entry['period_human'])}: ожидание
         <b>{_fmt(entry['expected_l'], 'л')}</b> при факте
         <b>{_fmt(entry['actual_l'], 'л')}</b> →
         {_fmt(abs(entry['saved_l']), 'л')} / <b>{_fmt(abs(entry['saved_kzt']))} ₸</b> — {sign}.</p>
      <p>Компоненты: {comp_bits}.</p>
      <p class="muted">Baseline заморожен {escape(str(b.get('frozen_at', ''))[:10])}:
         холостой ход {rates.get('idle_share', 0) * 100:.0f}%, движение
         {rates.get('moving_l_per_100km', 0):.1f} л/100 км
         (сезонная поправка ×{entry.get('season_factor', 1.0):.2f}).
         Методика: ожидание = ставки baseline × фактический пробег и моточасы
         периода — «ездили меньше» экономией не считается.</p>
    </section>"""


def _add_signals_slide(prs: Presentation, report: FleetReport) -> None:
    """Сигналы руководителю + топ проблемных ТС по баллу."""
    lines = list(report.alerts)
    top = [c for c in report.scorecard if c["score"] > 0][:6]
    if top:
        lines.append("— Топ внимания (балл проблемности):")
        lines += [f"   {c['name']}: {c['score']:.0f} — {c['reasons']}" for c in top]
    if not lines:
        lines = ["Критических сигналов за период не зафиксировано."]
    _add_bullets_slide(prs, report, "Сигналы и приоритеты", lines,
                       intro="Автоматические сигналы по порогам и рейтинг ТС по «проблемности».")


def _add_whatif_slide(prs: Presentation, report: FleetReport) -> None:
    """Потенциал экономии («что если») + бенчмарк по клиентам."""
    lines = []
    for s in report.whatif:
        if s["saved_kzt"]:
            lines.append(f"Сократить простой на {s['cut'] * 100:.0f}% → "
                         f"экономия {_fmt(s['saved_kzt'])} ₸ ({_fmt(s['saved_l'], 'л', 1)}).")
    b = report.benchmark
    if b.get("metrics"):
        lines.append(f"— Сравнение со средним по {b['peers']} паркам:")
        for m in b["metrics"].values():
            mark = "лучше" if m["better"] else "хуже"
            lines.append(f"   {m['label']}: у вас {m['mine']}, в среднем {m['peers_avg']} — {mark}.")
    if lines:
        _add_bullets_slide(prs, report, "Потенциал экономии", lines,
                           intro="Оценка экономии при сокращении простоя и сравнение с другими парками.")


def _add_money_story_slides(prs: Presentation, report: FleetReport) -> None:
    """Слайды-нарратив «Главное»: боль + куда уходят деньги (с методом) и
    «Что делать». Открывают отчёт; без цены ГСМ не строятся."""
    s = _money_story(report)
    if not s:
        return

    # Слайд: боль + таблица утечек с методом расчёта.
    slide = _blank_slide(prs)
    _add_slide_header(slide, "Главное: куда уходят деньги")
    pain = (f"{_fmt(s['idle_cost'])} ₸ сожжено на холостом ходу — "
            f"{s['idle_share'] * 100:.0f}% всех денег на топливо. Двигатель "
            f"работал стоя {_fmt(s['idle_hours'], 'ч')}; в пересчёте на год — "
            f"≈ {_fmt(s['idle_annual'])} ₸ в стоянку.")
    _add_textbox(slide, MARGIN, Emu(1430000), Emu(SLIDE_W - 2 * MARGIN), Emu(620000),
                 pain, size=Pt(15), color=COLOR_ACCENT, bold=True)
    headers = ["Куда уходят деньги", "За период, ₸", "Как посчитано"]
    rows = [[f"{lk['label']} ({lk['tag']})", _fmt(lk["kzt"]), lk["method"]]
            for lk in s["leaks"]]
    _place_table(slide, headers, rows, top=2120000,
                 col_widths=[0.26, 0.18, 0.56],
                 foot="Срезы разные и не складываются в один итог. "
                      "«факт» — телеметрия; «нормы» — к заданным нормам; "
                      "«≈ оценка» — отраслевой коэффициент.")
    _add_footer(slide, report)

    # Слайд: что делать.
    if s["actions"]:
        lines = [f"{a['title']}. {a['detail']}" for a in s["actions"]]
        _add_bullets_slide(prs, report, "Что делать", lines,
                           intro="Приоритетные действия с привязкой к деньгам "
                                 "и конкретным ТС.")


def _place_table(slide, headers, rows, *, top, col_widths=None, foot=None):
    """Таблица на ГОТОВОМ слайде (в отличие от _add_table_slide, который создаёт
    свой слайд). Шапка — синяя, тело — зебра; ширины колонок — доли ширины."""
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl_w = SLIDE_W - 2 * MARGIN
    graphic = slide.shapes.add_table(n_rows, n_cols, MARGIN, Emu(top),
                                     Emu(tbl_w), Emu(min(3800000, 360000 * n_rows)))
    table = graphic.table
    # Ширины колонок задаём ДО измерения высот (от них зависит перенос).
    if col_widths:
        for c, frac in enumerate(col_widths):
            table.columns[c].width = Emu(int(tbl_w * frac))
    # Последняя колонка — «как посчитано» (длинный текст), ей перенос нужен;
    # остальные держим в одну строку. Первую колонку клипуем.
    wrap_col = n_cols - 1
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        _style_run(r, size=SIZE_SMALL, color=COLOR_WHITE, bold=True)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_SOFT if ri % 2 == 0 else COLOR_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if c != wrap_col:
                cell.text_frame.word_wrap = False
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c != 1 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = _clip(val) if c == 0 else val
            _style_run(r, size=SIZE_TABLE_BODY, color=COLOR_TEXT)
    # Позиционируем по реальным высотам (учёт переноса колонки «как посчитано»).
    tbl_top, tbl_h = _fit_table(graphic, table, top0=top, footnote=bool(foot),
                                body_pt=SIZE_TABLE_BODY.pt)
    if foot:
        foot_top = min(tbl_top + tbl_h + 260000, FOOTER_TOP - 360000)
        _add_textbox(slide, MARGIN, Emu(foot_top),
                     Emu(tbl_w), Emu(330000), foot,
                     size=SIZE_FOOTER, color=COLOR_MUTED)


def _add_economics_slide(prs: Presentation, report: FleetReport) -> None:
    """Слайд «Экономический эффект»: корзины денег + COI (STRATEGY §4.1).

    Таблица корзин: «уже теряете за период» / «потенциал возврата», итог —
    адресуемые потери в ₸/месяц (cost of inaction). Оценочные корзины
    помечаются «≈». Без цены топлива слайд не строится.
    """
    from omnicomm_report import economics as econ_mod

    eco = econ_mod.build_economics(report)
    if not eco.buckets or eco.total_existing_kzt <= 0:
        return

    headers = ["Корзина", "Потери за период, ₸", "Потенциал возврата, ₸"]
    rows: list[list[str]] = []
    for b in eco.buckets:
        pre = "≈ " if b.is_estimate else ""   # маркер оценки — на ЗНАЧЕНИИ
        rows.append([
            b.label,
            (pre + _fmt(b.existing_kzt)) if b.existing_kzt > 0 else "—",
            (pre + _fmt(b.potential_kzt)) if b.potential_kzt > 0 else "—",
        ])
    rows.append([
        "ИТОГО", _fmt(eco.total_existing_kzt), _fmt(eco.total_potential_kzt),
    ])

    # Честная подача: строки «≈» — отраслевая ОЦЕНКА по коэффициентам (не факт),
    # ₸/мес — ЭКСТРАПОЛЯЦИЯ периода (×30/12), а не измеренный результат.
    legend = ("«≈» — оценка по отраслевым коэффициентам (не измерение); "
              "потенциал — доведение худших ТС до медианы парка. "
              "Измеренный результат программы — на слайде «Счётчик экономии».")
    # Экстраполяция в ₸/мес только на периоде ≥ 7 дней — короткий период
    # умножает шум (однодневный отчёт × 30 даёт абсурдные суммы).
    if eco.period_days >= 7:
        foot = (f"Оценка потерь без программы (экстраполяция периода): "
                f"≈ {_fmt(eco.coi_monthly_kzt)} ₸/мес "
                f"(≈ {_fmt(eco.coi_annual_kzt)} ₸/год). " + legend)
    else:
        foot = ("Суммы — за период отчёта (без экстраполяции на короткой "
                "выборке). " + legend)
    _add_table_slide(prs, report, "Экономический эффект: оценка потерь и потенциала",
                     headers, rows, footnote=foot)

    if eco.worst_vehicles:
        lines = [f"{name}: потенциал {_fmt(kzt)} ₸ за период"
                 for name, kzt in eco.worst_vehicles]
        _add_bullets_slide(
            prs, report, "Первоочередные ТС программы экономии", lines,
            intro=("Холостой ход выше медианы парка "
                   f"({eco.median_idle_share * 100:.0f}%). Значения — для "
                   "приоритизации проверки, не вывод о нарушениях."),
        )


def _add_savings_slide(prs: Presentation, report: FleetReport,
                       charts: dict[str, str]) -> None:
    """Слайд «Счётчик подтверждённой экономии» (Ф2, STRATEGY §4.2).

    Накопленный итог против замороженного baseline + результат периода
    по компонентам. Перерасход к эталону показывается честно. Без
    baseline слайд не строится.
    """
    s = report.savings or {}
    entry = s.get("period")
    if not entry:
        return

    cum_kzt = s.get("cumulative_kzt", 0)
    cum_l = s.get("cumulative_l", 0)
    n = s.get("entries_count", 0)
    b = s.get("baseline", {})
    rates = b.get("rates", {})
    src_periods = int(b.get("source_periods", 0) or 0)

    sign = "экономия" if entry["saved_l"] >= 0 else "перерасход к эталону"
    headline = (f"Накоплено за {n} пер.: "
                f"{_fmt(abs(cum_kzt))} ₸ ({_fmt(abs(cum_l), 'л', 0)}) — "
                + ("экономия" if cum_kzt >= 0 else "перерасход к эталону")
                + " против baseline")
    comp = entry.get("components", {})
    comp_bits = [f"{label} {comp[key].get('saved_l', 0):+,.0f} л"
                 for key, label in (("idle", "холостой ход"),
                                    ("moving", "движение"),
                                    ("stationary", "спецтехника"))
                 if key in comp]
    extra = [
        f"Период {entry['period_human']}: ожидание {_fmt(entry['expected_l'], 'л')} "
        f"при факте {_fmt(entry['actual_l'], 'л')} → "
        f"{_fmt(abs(entry['saved_l']), 'л')} / {_fmt(abs(entry['saved_kzt']))} ₸ — {sign}.",
        ("Компоненты: " + "; ".join(comp_bits) + ".").replace(",", " "),
        (f"Baseline заморожен {str(b.get('frozen_at', ''))[:10]} из "
         f"{src_periods} пер.: холостой ход "
         f"{rates.get('idle_share', 0) * 100:.0f}%, движение "
         f"{rates.get('moving_l_per_100km', 0):.1f} л/100 км "
         f"(сезонная поправка ×{entry.get('season_factor', 1.0):.2f})."),
        "Методика IPMVP-lite: ожидание = ставки baseline × фактический пробег и "
        "моточасы периода — «ездили меньше» экономией не считается.",
    ]
    # Честная пометка прочности эталона: 1 период baseline = предварительно.
    if src_periods < 2:
        extra.append("Baseline заморожен из одного периода — оценка "
                     "предварительная; точность растёт с накоплением истории.")
    _add_image_slide(
        prs, report, "Счётчик экономии против baseline",
        charts.get("savings"), headline, extra_lines=extra,
    )


# --- Публичные функции -------------------------------------------------------

def build_pptx(report: FleetReport, charts: dict[str, str], out_path: str) -> str:
    """Собрать клиентский .pptx по 8 слайдам ТЗ §8 и сохранить в out_path.

    `charts` — словарь путей к PNG от charts.build_charts с ключами
    'mileage','fuel_per_100km','fuel_idle','speeding'. Отсутствующий или
    несуществующий путь → слайд с плашкой «График недоступен» (не падаем).
    Возвращает абсолютный путь к созданному файлу.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    kpi = report.kpi

    # 1. Титул.
    _add_title_slide(prs, report)

    # 2. Ключевые выводы для руководства.
    summary = [
        f"Общий пробег парка: {_fmt(kpi.total_mileage_km, 'км')}.",
        f"Общий расход топлива: {_fmt(kpi.total_fuel_l, 'л')} "
        f"(средний по мобильным {_fmt(kpi.mobile_fuel_per_100km, 'л/100км', 1)}; "
        f"спецтехника — л/моточас в разделе использования).",
        f"Суммарные моточасы: {_fmt(kpi.total_engine_hours, 'ч', 1)} "
        f"(холостой ход {kpi.idle_hours_share * 100:.0f}%).",
    ]
    if kpi.weighted_fuel_per_motorhour > 0:
        summary.append(
            f"Средний расход на моточас работы: "
            f"{_fmt(kpi.weighted_fuel_per_motorhour, 'л/ч', 1)}."
        )
    if kpi.total_fuel_cost > 0:
        summary.append(
            f"Стоимость топлива: {_fmt(kpi.total_fuel_cost)} ₸ "
            f"(на простоях {_fmt(kpi.idle_fuel_cost)} ₸)."
        )
    trend_lines = _trend_lines(report)
    # Для руководства держим ёмкий список (≤6) — не «стена текста».
    # Полная детализация раскрывается на тематических слайдах ниже.
    key_points = summary + trend_lines + list(report.conclusions)
    _add_bullets_slide(
        prs, report, "Ключевые выводы для руководства", key_points[:6],
    )

    # 2-нарратив. Главное: боль → куда уходят деньги → что делать.
    _add_money_story_slides(prs, report)

    # 2a. Экономический эффект: корзины денег + COI (docs/STRATEGY.md §4.1).
    _add_economics_slide(prs, report)

    # 2b. Счётчик подтверждённой экономии против baseline (Ф2, STRATEGY §4.2).
    _add_savings_slide(prs, report, charts)

    # 3. Профиль автопарка (таблица, возможно несколько слайдов).
    _add_fleet_profile_slides(prs, report)

    # 3a. Классификация техники по типам (специфика анализа).
    _add_classification_slide(prs, report)

    # 3b. Структура парка: подвижная (л/100км) vs спецтехника (л/моточас).
    _add_image_slide(
        prs, report, "Структура парка: подвижная и спецтехника",
        charts.get("fleet_class"),
        f"Подвижная техника — {kpi.mobile_count} ТС (метрика л/100 км); "
        f"спецтехника — {kpi.stationary_count} ТС (метрика л/моточас).",
        extra_lines=[
            "Для неподвижной техники расход на 100 км не считается — "
            "при околонулевом пробеге он бессмыслен; её расход меряется в "
            "литрах на моточас работы.",
        ],
    )

    # 4. Пробег и распределение нагрузки.
    _add_image_slide(
        prs, report, "Пробег и распределение нагрузки",
        charts.get("mileage"),
        "Распределение пробега показывает загрузку парка по ТС.",
        extra_lines=[
            f"Максимальный вклад: {kpi.top_fuel_vehicle or '—'}."
            if kpi.top_fuel_vehicle else "Нагрузка распределена по парку.",
        ],
    )

    # 5. Топливная эффективность — подвижная техника (л/100 км).
    high_fuel = _high_fuel_vehicles(report)
    _add_image_slide(
        prs, report, "Топливная эффективность — подвижная техника",
        charts.get("fuel_per_100km"),
        "Подвижные ТС с наибольшим относительным расходом на 100 км:",
        extra_lines=(high_fuel or ["Данные о расходе на 100 км недоступны."])
        + [TXT_NORMS_DISCLAIMER],
    )

    # 5-bis. Топливная эффективность спецтехники (л/моточас) — если она есть.
    if _has_stationary_equipment(report):
        spec_lines = _spec_fuel_per_mh_lines(report)
        _add_image_slide(
            prs, report, "Топливная эффективность — спецтехника",
            charts.get("fuel_per_mh"),
            "Неподвижная спецтехника измеряется в литрах на моточас работы:",
            extra_lines=(spec_lines or ["Данные о расходе на моточас недоступны."])
            + [TXT_NORMS_DISCLAIMER],
        )

    # 5a. Использование парка (моточасы: движение vs холостой ход).
    idle_top = analytics.rank_idle(report.vehicles, top_n=5)
    util_lines = [
        f"«{n}»: холостой ход {s * 100:.0f}% "
        f"({_fmt(h, 'ч', 1) if h is not None else '—'})"
        for n, s, h in idle_top if s
    ]
    _add_image_slide(
        prs, report, "Использование парка",
        charts.get("utilization"),
        f"Из {_fmt(kpi.total_engine_hours, 'моточасов', 1)} в движении — "
        f"{_fmt(kpi.movement_hours, 'ч', 1)}, на холостом ходу — "
        f"{_fmt(kpi.total_idle_hours, 'ч', 1)} ({kpi.idle_hours_share * 100:.0f}%).",
        extra_lines=(["Наибольшая доля холостого хода:"] + util_lines
                     if util_lines else
                     ["Данные о холостом ходе по ТС недоступны."]),
    )

    # 6. Расход без движения.
    idle_share_pct = kpi.fuel_idle_share * 100
    _add_image_slide(
        prs, report, "Расход топлива без движения",
        charts.get("fuel_idle"),
        f"Объём топлива без движения: {_fmt(kpi.fuel_idle_l, 'л', 1)} "
        f"({idle_share_pct:.1f}% от общего расхода).",
        extra_lines=[
            "Рекомендация: контролировать длительные простои с работающим "
            "двигателем, согласовать регламент прогрева и стоянок.",
        ],
    )

    # 6a. Финансовая оценка (₸) — только при заданной цене топлива.
    if kpi.total_fuel_cost > 0:
        # Куда уходят деньги: топливо в движении vs на простое (наглядный стэк).
        _add_image_slide(
            prs, report, "Куда уходят деньги (₸)",
            charts.get("money"),
            f"Стоимость топлива за период: {_fmt(kpi.total_fuel_cost)} ₸ "
            f"по цене {_fmt(kpi.fuel_price_kzt)} ₸/л.",
            extra_lines=_financial_lines(report) + [TXT_NORMS_DISCLAIMER],
        )
        # Детализация стоимости по ТС.
        _add_image_slide(
            prs, report, "Стоимость топлива по ТС (₸)",
            charts.get("cost"),
            "Распределение затрат на топливо по машинам за период.",
        )

    # 6b. Отклонение от норм (перерасход/экономия) — если нормы заданы.
    if kpi.vehicles_with_norm > 0:
        _add_image_slide(
            prs, report, "Перерасход и экономия к норме (₸)",
            charts.get("overrun"),
            f"Нормы заданы у {kpi.vehicles_with_norm} ТС. "
            f"Перерасход: {_fmt(kpi.total_overrun_cost)} ₸; "
            f"экономия: {_fmt(kpi.total_economy_cost)} ₸.",
            extra_lines=[
                "Янтарь — перерасход к норме, синий — экономия. "
                "Норма учитывает сезонные и эксплуатационные коэффициенты.",
            ],
        )

    # 7. Операционные отклонения (только REVIEW, без обвинений).
    # Пустой слайд не рендерим: нужен либо график, либо аномалии, либо
    # ненулевое превышение скорости — иначе слайд не несёт смысла.
    anomalies = _review_anomalies(report)
    speeding_pct = kpi.speeding_mileage_share * 100
    speeding_chart = charts.get("speeding")
    has_speeding_chart = bool(speeding_chart and os.path.exists(speeding_chart))
    if has_speeding_chart or anomalies or speeding_pct > 0:
        _add_image_slide(
            prs, report, "Операционные отклонения",
            speeding_chart,
            f"Пробег с превышением скорости: {speeding_pct:.1f}% "
            f"(макс. зафиксированная скорость {_fmt(kpi.max_speed_kmh, 'км/ч', 0)}).",
            extra_lines=(anomalies or ["Отклонений, требующих проверки, не выявлено."])
            + ["Перечисленные значения требуют проверки корректности данных, "
               "а не являются выводом о нарушениях."],
        )

    # 7a. Работа на погрузке (мусоровозы/спецтехника) — если есть сигнал.
    if _loading_relevant(report):
        _add_loading_slides(prs, report, charts)

    # 7b. Перерасход / экономия по нормам — если нормы заданы.
    if report.kpi.vehicles_with_norm > 0:
        _add_norms_slide(prs, report)

    # 7c. Сигналы/приоритеты и потенциал экономии.
    _add_signals_slide(prs, report)
    _add_whatif_slide(prs, report)

    # 8. Рекомендации и план действий (динамические — из фактических KPI).
    _add_bullets_slide(
        prs, report, "Рекомендации и план действий",
        report.recommendations or [
            "Утвердить нормы расхода по типам ТС.",
            "Контролировать простои с работающим двигателем.",
            "Ввести регулярную отчётность для отслеживания динамики.",
        ],
    )

    abspath = os.path.abspath(out_path)
    os.makedirs(os.path.dirname(abspath) or ".", exist_ok=True)
    prs.save(abspath)
    logger.info("PPTX сохранён: %s (%d слайдов)", abspath, len(prs.slides._sldIdLst))
    return abspath


def export_xlsx(report: FleetReport, out_path: str) -> str:
    """Выгрузить очищенные/рассчитанные данные парка в .xlsx (ТЗ §10).

    Включает рассчитанный расход на 100 км (fuel_per_100km_calc) и пометку
    «нет данных». Колонка сливов НЕ выгружается (инвариант ТЗ §7).
    Лист KPI — агрегаты по парку. Возвращает абсолютный путь.
    """
    rows = []
    for v in report.vehicles:
        rows.append({
            "ТС": v.name,
            "Тип": vehicle_types.label(v.vehicle_type),
            "Марка": v.brand or "",
            "Модель": v.model or "",
            "Год": v.year or "",
            "Госномер": v.reg_number or "",
            "Двигатель": v.engine_model or "",
            "Бак, л": v.tank_capacity_l or "",
            "Норма л/100км": v.norm_l_per_100km or "",
            "Норма л/мч": v.norm_l_per_mh or "",
            "Отклонение, л": v.overrun_l if v.overrun_l is not None else "",
            "Группа": v.group or "",
            "Есть данные": "да" if v.has_data else TXT_NO_DATA,
            "Причина отсутствия": v.no_data_reason or "",
            "Пробег, км": v.mileage_km,
            "Расход, л": v.fuel_l,
            "Расход на 100 км (расч.)": v.fuel_per_100km_calc,
            "Расход на моточас (расч.)": v.fuel_per_motorhour,
            "Моточасы": v.engine_hours,
            "Холостой ход, ч": v.engine_idle_hours,
            "Топливо без движения, л": v.fuel_idle_l,
            "Макс. скорость, км/ч": v.max_speed_kmh,
            "Превышения, шт": v.speeding_count,
            "Пробег с превышением, км": v.speeding_mileage_km,
            "Аномалий (требуют проверки)": sum(
                1 for a in v.anomalies if a.severity == Severity.REVIEW
            ),
        })
    df_vehicles = pd.DataFrame(rows)

    kpi = report.kpi
    df_kpi = pd.DataFrame([
        ("Клиент", report.client_name),
        ("Период", report.period.human()),
        ("Источник данных", report.source),
        ("ТС всего", kpi.vehicles_total),
        ("ТС с данными", kpi.vehicles_with_data),
        ("Общий пробег, км", round(kpi.total_mileage_km, 1)),
        ("Общий расход, л", round(kpi.total_fuel_l, 1)),
        ("Средний расход, л/100км", round(kpi.weighted_fuel_per_100km, 1)),
        ("Моточасы всего", round(kpi.total_engine_hours, 1)),
        ("Топливо без движения, л", round(kpi.fuel_idle_l, 1)),
        ("Доля расхода без движения, %", round(kpi.fuel_idle_share * 100, 1)),
        ("Холостой ход, моточасов", round(kpi.total_idle_hours, 1)),
        ("Доля холостого хода, %", round(kpi.idle_hours_share * 100, 1)),
        ("Доля пробега с превышением, %", round(kpi.speeding_mileage_share * 100, 1)),
        ("Макс. скорость, км/ч", round(kpi.max_speed_kmh, 1)),
        ("Цена топлива, ₸/л", round(kpi.fuel_price_kzt, 1)),
        ("Стоимость топлива, ₸", round(kpi.total_fuel_cost, 0)),
        ("Топливо на простоях, ₸", round(kpi.idle_fuel_cost, 0)),
        ("Потенциальная экономия, ₸", round(kpi.potential_savings, 0)),
    ], columns=["Показатель", "Значение"])

    abspath = os.path.abspath(out_path)
    os.makedirs(os.path.dirname(abspath) or ".", exist_ok=True)
    with pd.ExcelWriter(abspath, engine="openpyxl") as writer:
        df_kpi.to_excel(writer, sheet_name="KPI", index=False)
        df_vehicles.to_excel(writer, sheet_name="ТС", index=False)
    logger.info("XLSX сохранён: %s (%d ТС)", abspath, len(rows))
    return abspath


def export_pdf(pptx_path: str, out_path: str) -> Optional[str]:
    """Сконвертировать .pptx → .pdf через LibreOffice CLI (ТЗ §10, опционально).

    PDF-экспорт необязателен: если LibreOffice (soffice) не установлен или
    конвертация не удалась — возвращаем None и логируем, НЕ падаем.
    Возвращает абсолютный путь к .pdf при успехе.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        logger.warning("PDF-экспорт пропущен: LibreOffice (soffice) не найден в PATH.")
        return None

    src = os.path.abspath(pptx_path)
    if not os.path.exists(src):
        logger.warning("PDF-экспорт пропущен: исходный .pptx не найден: %s", src)
        return None

    abspath = os.path.abspath(out_path)
    out_dir = os.path.dirname(abspath) or "."
    os.makedirs(out_dir, exist_ok=True)
    try:
        # LibreOffice headless кладёт PDF в out_dir под именем исходника.
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, src],
            check=True, capture_output=True, timeout=120,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as exc:
        logger.warning("PDF-экспорт не удался (%s), отчёт остаётся в .pptx.", exc)
        return None

    produced = os.path.join(out_dir, os.path.splitext(os.path.basename(src))[0] + ".pdf")
    if not os.path.exists(produced):
        logger.warning("PDF-экспорт: LibreOffice не создал ожидаемый файл %s.", produced)
        return None
    if os.path.abspath(produced) != abspath:
        shutil.move(produced, abspath)
    logger.info("PDF сохранён: %s", abspath)
    return abspath


# --- HTML-отчёт (ТЗ §8, альтернатива PDF) ------------------------------------
# Самодостаточный HTML: графики встроены base64, шрифты системные. Открывается
# в любом браузере; печать в PDF — Ctrl+P → «Сохранить как PDF». Цвета и
# формулировки те же, что в .pptx (палитра ниже дублирует RGBColor в hex).

HTML_PRIMARY = "#2F5C8F"
HTML_TEXT = "#2B2B2B"
HTML_ACCENT = "#C8893F"
HTML_BG_SOFT = "#F4F6F9"
HTML_MUTED = "#8A8A8A"
HTML_BORDER = "#E2E6EC"
# Serif-дисплей для крупных чисел/заголовков (editorial financial dossier).
SERIF_STACK = "Georgia, 'Iowan Old Style', 'Times New Roman', serif"


def _img_data_uri(path: Optional[str]) -> Optional[str]:
    """PNG-файл → data:URI для встраивания в HTML; None если файла нет."""
    if not path or not os.path.exists(path):
        return None
    with open(path, "rb") as fh:
        b64 = base64.b64encode(fh.read()).decode("ascii")
    return f"data:image/png;base64,{b64}"


def _inline_svg(png_path: Optional[str]) -> Optional[str]:
    """SVG-сосед PNG-файла → инлайн-разметка (чёткий вектор, печать в PDF).

    matplotlib пишет .svg рядом с .png (см. charts._save). Инлайним вектор:
    срезаем XML-пролог/DOCTYPE и делаем его адаптивным (width:100%).
    """
    if not png_path:
        return None
    svg_path = os.path.splitext(png_path)[0] + ".svg"
    if not os.path.exists(svg_path):
        return None
    try:
        with open(svg_path, encoding="utf-8") as fh:
            svg = fh.read()
    except OSError:
        return None
    idx = svg.find("<svg")
    if idx < 0:
        return None
    svg = svg[idx:]  # отрезать <?xml ...?> и <!DOCTYPE ...>
    # Адаптивный размер: viewBox от matplotlib сохраняется, фикс. pt-размеры убираем.
    svg = svg.replace("<svg ", '<svg style="width:100%;height:auto;display:block" ', 1)
    return svg


def _html_figure(charts: dict[str, str], key: str, alt: str) -> str:
    """Блок с графиком: инлайн-SVG (вектор) при наличии, иначе base64-PNG."""
    svg = _inline_svg(charts.get(key))
    if svg:
        return f'<figure role="img" aria-label="{escape(alt)}">{svg}</figure>'
    uri = _img_data_uri(charts.get(key))
    if uri:
        return f'<figure><img src="{uri}" alt="{escape(alt)}"></figure>'
    return f'<div class="chart-missing">{TXT_CHART_UNAVAILABLE}</div>'


def _kpi_delta(report: FleetReport, metric: Optional[str],
              worse_if_up: Optional[bool]) -> str:
    """Бэйдж дельты период-к-периоду для KPI-плитки: ▲/▼ X% с цветовым статусом.

    Цвет: зелёный — изменение «к лучшему», янтарь — «к худшему», серый — нейтрально
    (без алармистского красного: инвариант «без обвинительных формулировок»).
    """
    if not metric:
        return ""
    delta = (report.trends or {}).get(metric)
    if delta is None or delta == 0:
        return ""
    up = delta > 0
    arrow = "▲" if up else "▼"
    if worse_if_up is None:
        cls = "d-flat"                       # нейтральная метрика (пробег)
    else:
        better = (not up) if worse_if_up else up
        cls = "d-good" if better else "d-warn"
    return (f'<div class="card-delta {cls}">{arrow} {abs(delta):.0f}% '
            f'<span class="d-note">к прошлому периоду</span></div>')


def _html_exec_band(report: FleetReport) -> str:
    """«Экран для руководителя»: деньги и холостой ход одним взглядом (hero-SVG)."""
    kpi = report.kpi
    money = svg_charts.money_split(kpi, report)
    avg = None
    metrics = (report.benchmark or {}).get("metrics", {})
    if "idle_hours_share" in metrics:
        avg = metrics["idle_hours_share"].get("peers_avg")
    worst = analytics.rank_idle(report.vehicles, top_n=3)
    bullet = svg_charts.idle_bullet(kpi, worst, avg=avg)
    cards = []
    if money:
        cards.append(f'<div class="exec-card reveal">{money}</div>')
    if bullet:
        cards.append(f'<div class="exec-card reveal">{bullet}</div>')
    if not cards:
        return ""
    return f"""
    <section class="exec">
      <div class="exec-head">Сводка для руководителя</div>
      <div class="exec-grid">{''.join(cards)}</div>
    </section>
    """


def _html_list(items: list[str]) -> str:
    return "<ul>" + "".join(f"<li>{escape(s)}</li>" for s in items) + "</ul>"


def _html_profile_table(report: FleetReport) -> str:
    """Таблица профиля парка — без столбца «сливы» (бизнес-инвариант)."""
    head = (
        "<tr><th>ТС</th><th>Пробег, км</th><th>Расход, л</th><th>Расход уд.</th>"
        "<th>Моточасы</th><th>Простой, ч</th><th>Макс. V, км/ч</th></tr>"
    )
    rows: list[str] = []
    for v in report.vehicles:
        if not v.has_data:
            rows.append(
                f'<tr class="no-data"><td>{escape(v.name)}</td>'
                f'<td colspan="6">{TXT_NO_DATA}'
                + (f" ({escape(v.no_data_reason)})" if v.no_data_reason else "")
                + "</td></tr>"
            )
            continue
        prof = vehicle_types.profile(v.vehicle_type)
        rate = (f"{_fmt(v.fuel_per_motorhour, digits=1)} л/мч"
                if prof.primary_metric == "l_per_mh"
                else f"{_fmt(v.fuel_per_100km_calc, digits=1)} л/100км")
        cells = [
            escape(v.name),
            _fmt(v.mileage_km),
            _fmt(v.fuel_l),
            rate,
            _fmt(v.engine_hours, digits=1),
            _fmt(v.engine_idle_hours, digits=1),
            _fmt(v.max_speed_kmh),
        ]
        rows.append("<tr>" + "".join(f"<td>{c}</td>" for c in cells) + "</tr>")
    return f"<table class='profile'>{head}{''.join(rows)}</table>"


def build_html(report: FleetReport, charts: dict[str, str], out_path: str) -> str:
    """Собрать клиентский HTML-отчёт (зеркало 8 разделов .pptx) и сохранить.

    Самодостаточный файл: графики встроены как base64-PNG, внешних зависимостей
    нет. Возвращает абсолютный путь. `charts` — те же ключи, что у build_pptx.
    """
    kpi = report.kpi
    generated = report.generated_at.strftime("%d.%m.%Y") if report.generated_at else ""

    # KPI-плитки в стиле аналитического дашборда: значение + дельта-тренд (▲/▼)
    # с цветовым статусом. (val, label, trend_metric|None, worse_if_up, status)
    kpi_cards = [
        (_fmt(kpi.total_mileage_km, "км"), "Общий пробег парка", "total_mileage_km", None, ""),
        (_fmt(kpi.total_fuel_l, "л"), "Общий расход топлива", "total_fuel_l", True, ""),
        (_fmt(kpi.weighted_fuel_per_100km, "л/100км", 1), "Средневзвешенный расход",
         "weighted_fuel_per_100km", True, ""),
        (f"{kpi.idle_hours_share * 100:.0f}%", "Доля холостого хода", "idle_hours_share",
         True, "warn" if kpi.idle_hours_share > 0.30 else ""),
    ]
    if kpi.weighted_fuel_per_motorhour > 0:
        kpi_cards.append((_fmt(kpi.weighted_fuel_per_motorhour, "л/ч", 1),
                          "Расход на моточас", "weighted_fuel_per_motorhour",
                          True, ""))
    if kpi.utilization_fund > 0:
        kpi_cards.append((f"{kpi.utilization_fund * 100:.0f}%",
                          f"Использование (фонд {kpi.time_fund_hours_per_day:g} ч/сут)",
                          None, None,
                          "good" if kpi.utilization_fund >= 0.7 else "warn"))
    elif kpi.utilization_calendar > 0:
        kpi_cards.append((f"{kpi.utilization_calendar * 100:.0f}%",
                          "Использование (календарное)", None, None, ""))
    if kpi.fuel_cost_per_m3 > 0:
        kpi_cards.append((f"{_fmt(kpi.fuel_cost_per_m3)} ₸/м³",
                          f"Топливо на 1 м³ вывоза ({_fmt(kpi.haul_volume_m3)} м³)",
                          None, None, ""))
    if kpi.total_fuel_cost > 0:
        kpi_cards.append((f"{_fmt(kpi.total_fuel_cost)} ₸", "Стоимость топлива",
                          "total_fuel_cost", True, ""))
        kpi_cards.append((f"{_fmt(kpi.potential_savings)} ₸", "Потенциальная экономия",
                          None, None, "good"))
    cards_html = "".join(
        f'<div class="card{" card-" + status if status else ""}">'
        f'<div class="card-val">{escape(val)}</div>'
        f'<div class="card-lbl">{escape(lbl)}</div>'
        f'{_kpi_delta(report, metric, worse_if_up)}</div>'
        for val, lbl, metric, worse_if_up, status in kpi_cards
    )

    summary = [
        f"Общий пробег парка: {_fmt(kpi.total_mileage_km, 'км')}.",
        f"Общий расход топлива: {_fmt(kpi.total_fuel_l, 'л')} "
        f"(средний по мобильным {_fmt(kpi.mobile_fuel_per_100km, 'л/100км', 1)}; "
        f"спецтехника — л/моточас в разделе использования).",
        f"Суммарные моточасы: {_fmt(kpi.total_engine_hours, 'ч', 1)} "
        f"(холостой ход {kpi.idle_hours_share * 100:.0f}%).",
    ]
    if kpi.total_fuel_cost > 0:
        summary.append(
            f"Стоимость топлива: {_fmt(kpi.total_fuel_cost)} ₸, "
            f"из них на простоях {_fmt(kpi.idle_fuel_cost)} ₸."
        )
    summary += _trend_lines(report) + list(report.conclusions)

    high_fuel = _high_fuel_vehicles(report)
    idle_top = analytics.rank_idle(report.vehicles, top_n=5)
    idle_lines = [f"{n}: холостой ход {s * 100:.0f}%" for n, s, _ in idle_top if s]
    anomalies = _review_anomalies(report)
    idle_pct = kpi.fuel_idle_share * 100
    speeding_pct = kpi.speeding_mileage_share * 100

    recommendations = report.recommendations or [
        "Утвердить нормы расхода по типам ТС.",
        "Контролировать простои с работающим двигателем.",
        "Ввести регулярную отчётность для отслеживания динамики.",
    ]

    # Секция использования парка (всегда), финансов (₸) — при заданной цене.
    util_section = f"""
    <section>
      <h2>Использование парка</h2>
      {_html_figure(charts, "utilization", "Моточасы: движение и холостой ход")}
      <p>Из {_fmt(kpi.total_engine_hours, 'моточасов', 1)} в движении —
         {_fmt(kpi.movement_hours, 'ч', 1)}, на холостом ходу —
         {_fmt(kpi.total_idle_hours, 'ч', 1)} ({kpi.idle_hours_share * 100:.0f}%).</p>
      {_html_list(idle_lines or ["Данные о холостом ходе по ТС недоступны."])}
    </section>
    """
    # Эффективность спецтехники (л/моточас) — отдельная секция, если она в парке.
    spec_fuel_section = ""
    if _has_stationary_equipment(report):
        spec_lines = _spec_fuel_per_mh_lines(report)
        spec_fuel_section = f"""
    <section>
      <h2>Топливная эффективность — спецтехника</h2>
      {_html_figure(charts, "fuel_per_mh", "Расход на моточас (спецтехника)")}
      <p>Неподвижная спецтехника измеряется в литрах на моточас работы:</p>
      {_html_list(spec_lines or ["Данные о расходе на моточас недоступны."])}
      <p class="disclaimer">{escape(TXT_NORMS_DISCLAIMER)}</p>
    </section>
    """

    money_section = ""
    if kpi.total_fuel_cost > 0:
        fin_items = "".join(f"<li>{escape(s)}</li>" for s in _financial_lines(report))
        money_svg = svg_charts.money_split(kpi, report)
        money_fig = (f'<figure role="img" aria-label="Куда уходят деньги">{money_svg}</figure>'
                     if money_svg else _html_figure(charts, "money", "Куда уходят деньги"))
        money_section = f"""
    <section>
      <h2>Финансовая оценка (₸)</h2>
      {money_fig}
      <p>Стоимость топлива за период: <b>{_fmt(kpi.total_fuel_cost)} ₸</b>
         по цене {_fmt(kpi.fuel_price_kzt)} ₸/л.</p>
      <ul>{fin_items}</ul>
      {_html_figure(charts, "cost", "Стоимость топлива по ТС")}
      <p class="disclaimer">{escape(TXT_NORMS_DISCLAIMER)}</p>
    </section>
    """

    loading_section = ""
    if _loading_relevant(report):
        mh_stat = _fleet_fuel_per_mh_stationary(report)
        rows = _loading_rows(report)
        trow = "".join(
            "<tr>" + "".join(f"<td>{escape(str(c))}</td>" for c in r) + "</tr>"
            for r in rows
        )
        table = (
            "<table class='profile'><tr><th>ТС</th><th>Погрузка, ч</th>"
            "<th>Полезно %</th><th>Топл., л</th><th>л/моточас</th><th>Источник</th></tr>"
            + trow + "</table>"
        ) if rows else ""
        # Журнал обслуживания площадок (только при геозонах). GPS-карта точек убрана.
        if _has_geozone_visits(report):
            gz_rows = _geozone_journal_rows(report)
            gz_tr = "".join("<tr>" + "".join(f"<td>{escape(str(c))}</td>" for c in r) + "</tr>"
                            for r in gz_rows)
            map_block = ("<h3>Журнал обслуживания площадок (геозоны)</h3>"
                         "<table class='profile'><tr><th>Площадка</th><th>Визитов</th>"
                         "<th>Время</th><th>Топливо</th><th>₸</th></tr>" + gz_tr + "</table>")
        else:
            map_block = ""
        pts = sum(len(v.loading_points) for v in report.vehicles)
        loading_section = f"""
    <section>
      <h2>Работа на погрузке</h2>
      {_html_figure(charts, "loading_split", "Погрузка и простой")}
      <p>Время погрузки: <b>{_fmt(kpi.total_loading_hours_sensor, 'ч', 1)}</b> по датчику (факт){
          f" + {_fmt(kpi.total_loading_hours_estimate, 'ч', 1)} по оценке (обороты/GPS-маршрут, ≈)" if kpi.total_loading_hours_estimate else ""}.
         Доля полезной работы из времени стоя: <b>{kpi.fleet_loading_utilization * 100:.0f}%</b>.</p>
      {f'<p>Расход на моточас работы стоя (подъём контейнеров): <b>{_fmt(mh_stat, "л/моточас", 1)}</b> — топливо у такой техники списывается по моточасам, а не только по пробегу.</p>' if mh_stat is not None else ''}
      {f'<p>Потери на непродуктивном простое: <b>{_fmt(kpi.total_unproductive_fuel_cost)} ₸</b> ({_fmt(kpi.total_unproductive_fuel_l, "л", 1)}).</p>' if kpi.total_unproductive_fuel_cost > 0 else ''}
      <p class="muted">Датчик надстройки настроен у {kpi.vehicles_with_loading_sensor} из {kpi.vehicles_with_data} ТС; «≈» — оценка по оборотам/GPS-маршруту (без датчика).</p>
      {table}
      {map_block}
      {f'<p class="muted">Обслужено остановок-обслуживаний: {pts} (оценка по GPS-маршруту).</p>' if pts and not _has_geozone_visits(report) else ''}
    </section>
    """

    exec_band = _html_exec_band(report)

    sections = f"""
    <section class="hero">
      <div class="hero-eyebrow">Аналитический отчёт по автопарку</div>
      <h1>{escape(report.client_name)}</h1>
      <div class="hero-meta">Период: {escape(report.period.human())}{
          f' · сформирован {generated}' if generated else ''}</div>
      <div class="kpis">{cards_html}</div>
    </section>

    {_html_money_story(report)}

    {exec_band}

    <section>
      <h2>Ключевые выводы для руководства</h2>
      {_html_list(summary)}
    </section>

    {_html_economics_section(report)}

    {_html_savings_section(report, charts)}

    {_html_signals_section(report)}

    <section>
      <h2>Профиль автопарка</h2>
      {_html_profile_table(report)}
    </section>

    {_html_classification_section(report)}

    <section>
      <h2>Структура парка: подвижная и спецтехника</h2>
      {_html_figure(charts, "fleet_class", "Структура парка по типу метрики")}
      <p>Подвижная техника — <b>{kpi.mobile_count} ТС</b> (метрика л/100 км);
         спецтехника — <b>{kpi.stationary_count} ТС</b> (метрика л/моточас).</p>
      <p class="muted">Для неподвижной техники расход на 100 км не считается — при
         околонулевом пробеге он бессмыслен; её расход меряется в литрах на моточас.</p>
    </section>

    <section>
      <h2>Пробег и распределение нагрузки</h2>
      {_html_figure(charts, "mileage", "Распределение пробега")}
      <p class="muted">Распределение пробега показывает загрузку парка по ТС.</p>
    </section>

    <section>
      <h2>Топливная эффективность — подвижная техника</h2>
      {_html_figure(charts, "fuel_per_100km", "Расход на 100 км (подвижная техника)")}
      <p>Подвижные ТС с наибольшим относительным расходом на 100 км:</p>
      {_html_list(high_fuel or ["Данные о расходе на 100 км недоступны."])}
      <p class="disclaimer">{escape(TXT_NORMS_DISCLAIMER)}</p>
    </section>
    {spec_fuel_section}

    {util_section}

    <section>
      <h2>Расход топлива без движения</h2>
      {_html_figure(charts, "fuel_idle", "Расход без движения")}
      <p>Объём топлива без движения: {_fmt(kpi.fuel_idle_l, 'л', 1)}
         ({idle_pct:.1f}% от общего расхода).</p>
      <p class="muted">Рекомендация: контролировать длительные простои с
         работающим двигателем, согласовать регламент прогрева и стоянок.</p>
    </section>

    {money_section}

    {loading_section}

    {_html_norms_section(report, charts)}

    <section>
      <h2>Операционные отклонения</h2>
      {_html_figure(charts, "speeding", "Превышения скорости")}
      <p>Пробег с превышением скорости: {speeding_pct:.1f}%
         (макс. зафиксированная скорость {_fmt(kpi.max_speed_kmh, 'км/ч')}).</p>
      {_html_list(anomalies or ["Отклонений, требующих проверки, не выявлено."])}
      <p class="disclaimer">Перечисленные значения требуют проверки корректности
         данных, а не являются выводом о нарушениях.</p>
    </section>

    <section>
      <h2>Рекомендации и план действий</h2>
      {_html_list(recommendations)}
    </section>
    """

    html = f"""<!DOCTYPE html>
<html lang="ru">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Отчёт по автопарку — {escape(report.client_name)}</title>
<style>
  :root {{ color-scheme: light; }}
  * {{ box-sizing: border-box; }}
  body {{ margin: 0; background: {HTML_BG_SOFT}; color: {HTML_TEXT};
    font-family: -apple-system, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
    line-height: 1.5; }}
  .wrap {{ max-width: 960px; margin: 0 auto; padding: 24px; }}
  section {{ background: #fff; border: 1px solid {HTML_BORDER}; border-radius: 12px;
    padding: 24px 28px; margin: 18px 0; }}
  .hero {{ border: none; position: relative; overflow: hidden; padding: 40px 34px 30px;
    background:
      radial-gradient(120% 140% at 88% -20%, #eaf1fa 0%, rgba(234,241,250,0) 60%),
      linear-gradient(180deg, #ffffff 0%, #fbfcfe 100%);
    box-shadow: 0 1px 2px rgba(20,40,70,.04), 0 18px 40px -26px rgba(20,40,70,.22); }}
  .hero::before {{ content: ""; position: absolute; left: 0; top: 0; bottom: 0; width: 5px;
    background: linear-gradient(180deg, {HTML_PRIMARY}, {HTML_ACCENT}); }}
  .hero-eyebrow {{ font-size: 12.5px; font-weight: 700; letter-spacing: .14em;
    text-transform: uppercase; color: {HTML_ACCENT}; }}
  h1 {{ color: {HTML_PRIMARY}; font-size: 34px; line-height: 1.1; margin: 8px 0 4px;
    font-family: {SERIF_STACK}; letter-spacing: -.01em; }}

  /* Одностраничник «Главное»: боль → куда уходят деньги → что делать */
  .story {{ border: none; padding: 30px 30px 34px;
    background: linear-gradient(180deg, #fff 0%, {HTML_BG_SOFT} 100%);
    box-shadow: inset 0 1px 0 #fff, 0 1px 2px rgba(20,40,70,.05); }}
  .story-eyebrow {{ font-size: 12.5px; font-weight: 700; letter-spacing: .14em;
    text-transform: uppercase; color: {HTML_MUTED}; margin-bottom: 14px; }}
  .pain {{ display: flex; align-items: baseline; gap: 20px; padding: 18px 22px;
    border-left: 5px solid {HTML_ACCENT}; background: #fff;
    border-radius: 0 12px 12px 0; box-shadow: 0 1px 2px rgba(20,40,70,.05);
    margin-bottom: 26px; flex-wrap: wrap; }}
  .pain-num {{ font-family: {SERIF_STACK}; font-size: 40px; font-weight: 700;
    color: {HTML_ACCENT}; line-height: 1; white-space: nowrap; }}
  .pain-txt {{ font-size: 16.5px; line-height: 1.5; flex: 1; min-width: 280px; }}
  .story-h {{ font-size: 13px; font-weight: 700; letter-spacing: .12em;
    text-transform: uppercase; color: {HTML_PRIMARY}; margin: 22px 0 12px; }}
  .leaks {{ display: flex; flex-direction: column; gap: 10px; }}
  .leak {{ display: flex; gap: 16px; align-items: flex-start; background: #fff;
    border: 1px solid {HTML_BORDER}; border-radius: 12px; padding: 14px 18px; }}
  .leak-rank {{ font-family: {SERIF_STACK}; font-size: 22px; font-weight: 700;
    color: {HTML_MUTED}; line-height: 1.2; min-width: 22px; }}
  .leak-body {{ flex: 1; }}
  .leak-top {{ display: flex; justify-content: space-between; align-items: baseline;
    gap: 12px; flex-wrap: wrap; }}
  .leak-label {{ font-size: 16px; font-weight: 700; color: {HTML_TEXT}; }}
  .leak-kzt {{ font-family: {SERIF_STACK}; font-size: 21px; font-weight: 700;
    color: {HTML_PRIMARY}; white-space: nowrap; }}
  .leak-method {{ font-size: 13.5px; color: {HTML_MUTED}; margin-top: 3px;
    line-height: 1.45; }}
  .leak-tag {{ font-size: 10.5px; font-weight: 700; letter-spacing: .03em;
    text-transform: uppercase; padding: 2px 7px; border-radius: 6px;
    margin-left: 9px; vertical-align: middle; }}
  .tag-fact {{ background: #e8f0f8; color: {HTML_PRIMARY}; }}
  .tag-norm {{ background: #f4ece0; color: {HTML_ACCENT}; }}
  .tag-est {{ background: #f0f0f0; color: {HTML_MUTED}; }}
  .story-note {{ font-size: 12.5px; color: {HTML_MUTED}; margin: 12px 2px 0;
    line-height: 1.45; }}
  .actions {{ margin: 4px 0 0; padding-left: 22px; }}
  .actions li {{ font-size: 15.5px; line-height: 1.55; margin-bottom: 10px; }}

  /* Экран для руководителя — «вау»-сводка */
  .exec {{ border: none; padding: 26px 28px 30px;
    background:
      radial-gradient(90% 160% at 0% 0%, #f3f7fc 0%, rgba(243,247,252,0) 55%),
      linear-gradient(180deg, #fdfefe 0%, {HTML_BG_SOFT} 100%);
    box-shadow: inset 0 1px 0 #fff, 0 1px 2px rgba(20,40,70,.04); }}
  .exec-head {{ font-size: 12.5px; font-weight: 700; letter-spacing: .14em;
    text-transform: uppercase; color: {HTML_MUTED}; margin-bottom: 16px; }}
  .exec-grid {{ display: grid; grid-template-columns: 1fr; gap: 16px; }}
  .exec-card {{ background: #fff; border: 1px solid {HTML_BORDER}; border-radius: 14px;
    padding: 22px 26px; box-shadow: 0 1px 2px rgba(20,40,70,.05),
    0 16px 36px -22px rgba(20,40,70,.28); }}

  /* Оркестрованная анимация на загрузке (print/reduced-motion отключают) */
  .reveal {{ animation: rise .7s cubic-bezier(.2,.7,.2,1) both; }}
  .exec-grid .exec-card:nth-child(2) {{ animation-delay: .12s; }}
  @keyframes rise {{ from {{ opacity: 0; transform: translateY(16px); }}
    to {{ opacity: 1; transform: none; }} }}
  svg .grow {{ transform-box: fill-box; transform-origin: left center;
    animation: grow 1s cubic-bezier(.2,.7,.2,1) both; }}
  svg .grow2 {{ animation-delay: .18s; }}
  @keyframes grow {{ from {{ transform: scaleX(0); }} to {{ transform: scaleX(1); }} }}
  @media (prefers-reduced-motion: reduce) {{
    .reveal, svg .grow {{ animation: none; opacity: 1; transform: none; }} }}
  h2 {{ color: {HTML_PRIMARY}; font-size: 20px; margin: 0 0 14px;
    border-bottom: 2px solid {HTML_BG_SOFT}; padding-bottom: 8px; }}
  .hero-sub {{ font-size: 18px; font-weight: 600; }}
  .hero-meta {{ color: {HTML_MUTED}; font-size: 14px; margin-bottom: 18px; }}
  .kpis {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 12px; }}
  .card {{ background: #fff; border: 1px solid {HTML_BORDER}; border-left: 3px solid {HTML_PRIMARY};
    border-radius: 10px; padding: 14px 16px;
    box-shadow: 0 1px 2px rgba(20,40,70,.04), 0 10px 24px -18px rgba(20,40,70,.25); }}
  .card.card-warn {{ border-left-color: {HTML_ACCENT}; }}
  .card.card-good {{ border-left-color: #3E7C5A; }}
  .card-val {{ color: {HTML_PRIMARY}; font-size: 24px; font-weight: 700;
    font-family: {SERIF_STACK}; letter-spacing: -.01em; font-variant-numeric: tabular-nums; }}
  .card-warn .card-val {{ color: {HTML_ACCENT}; }}
  .card-good .card-val {{ color: #3E7C5A; }}
  .card-lbl {{ color: {HTML_MUTED}; font-size: 12.5px; margin-top: 2px; }}
  .card-delta {{ font-size: 12px; font-weight: 700; margin-top: 8px;
    display: inline-flex; align-items: baseline; gap: 5px; }}
  .card-delta .d-note {{ font-weight: 400; color: {HTML_MUTED}; font-size: 11px; }}
  .d-good {{ color: #3E7C5A; }}
  .d-warn {{ color: {HTML_ACCENT}; }}
  .d-flat {{ color: {HTML_MUTED}; }}
  ul {{ margin: 0; padding-left: 22px; }}
  li {{ margin: 6px 0; }}
  figure {{ margin: 0 0 12px; text-align: center; }}
  img {{ max-width: 100%; height: auto; border: 1px solid {HTML_BORDER}; border-radius: 8px; }}
  .chart-missing {{ background: {HTML_BG_SOFT}; color: {HTML_MUTED};
    border: 1px dashed {HTML_BORDER}; border-radius: 8px; padding: 48px;
    text-align: center; margin-bottom: 12px; }}
  table.profile {{ width: 100%; border-collapse: collapse; font-size: 14px; }}
  table.profile th {{ background: {HTML_PRIMARY}; color: #fff; text-align: left;
    padding: 8px 10px; }}
  table.profile td {{ padding: 7px 10px; border-bottom: 1px solid {HTML_BORDER}; }}
  table.profile tr:nth-child(even) td {{ background: {HTML_BG_SOFT}; }}
  table.profile td:not(:first-child) {{ text-align: right;
    font-variant-numeric: tabular-nums; }}
  tr.no-data td {{ color: {HTML_MUTED}; font-style: italic; }}
  .muted {{ color: {HTML_MUTED}; }}
  .disclaimer {{ color: {HTML_MUTED}; font-size: 13px; font-style: italic; }}
  footer {{ color: {HTML_MUTED}; font-size: 12px; text-align: center; padding: 12px 0 4px; }}
  @media print {{
    body {{ background: #fff; }}
    .wrap {{ max-width: none; padding: 0; }}
    section {{ break-inside: avoid; border-radius: 0; box-shadow: none; }}
    .hero, .exec {{ box-shadow: none; }}
    .exec-card {{ box-shadow: none; }}
    .reveal, svg .grow {{ animation: none !important; opacity: 1 !important;
      transform: none !important; }}
  }}
  @media (max-width: 640px) {{ .kpis {{ grid-template-columns: repeat(2, 1fr); }} }}
</style>
</head>
<body>
<div class="wrap">
{sections}
<footer>Источник данных: Omnicomm Online{
    ' (REST API)' if report.source == 'api' else ' (выгрузка)'}.
  Значения «требуют проверки» — для приоритизации, не вывод о нарушениях.</footer>
</div>
</body>
</html>"""

    abspath = os.path.abspath(out_path)
    os.makedirs(os.path.dirname(abspath) or ".", exist_ok=True)
    with open(abspath, "w", encoding="utf-8") as fh:
        fh.write(html)
    logger.info("HTML сохранён: %s", abspath)
    return abspath
