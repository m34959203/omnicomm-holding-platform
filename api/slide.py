"""Выгрузка «галочкой»: выбранные разделы → ОДИН лист презентации (kb-15 §5).

Директор: «выбираешь галочку — нужная информация выгружается на одном листе
чётко и понятно, и сразу в виде презентации». Один слайд 16:9 из секций
УЖЕ ОТСКОУПЛЕННОГО снапшота (RBAC как у дашборда/Excel).

Бизнес-инварианты (.pptx): светлый официальный стиль; «сливы» НЕ выводятся;
без обвинительных формулировок; деньги крупно.
"""

from __future__ import annotations

import datetime as dt
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Палитра (тема Omnicomm Online, светлая)
INK = RGBColor(0x1B, 0x27, 0x33)
MUTED = RGBColor(0x5B, 0x6B, 0x80)
FAINT = RGBColor(0x8A, 0x98, 0xAC)
BLUE = RGBColor(0x1F, 0x6F, 0xD6)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
AMBER = RGBColor(0xD3, 0x9A, 0x1E)
RED = RGBColor(0xD4, 0x45, 0x3B)
LINE = RGBColor(0xE6, 0xE9, 0xEE)
PANEL = RGBColor(0xF7, 0xF9, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# Порядок и русские подписи секций (ключи — контракт с фронтом)
SECTIONS = {
    "fleet": "Парк",
    "economics": "Деньги",
    "speed": "Скоростной режим",
    "fuel": "Топливо",
    "quality": "Качество данных",
    "maint": "Контроль ТО",
    "tyres": "Шины",
}


def _ru(v, frac=0) -> str:
    """1234567.8 → «1 234 568» (неразрывные тонкие пробелы не нужны — pptx)."""
    if v is None:
        return "—"
    s = f"{float(v):,.{frac}f}".replace(",", " ").replace(".", ",")
    return s


def _mln(v) -> str:
    if v is None:
        return "—"
    v = float(v)
    if abs(v) >= 1e6:
        return f"{v / 1e6:.1f}".replace(".", ",") + " млн ₸"
    if abs(v) >= 1e3:
        return _ru(round(v / 1e3)) + " тыс ₸"
    return _ru(round(v)) + " ₸"


def _box(slide, x, y, w, h, *, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.045
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    return sh


def _text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, space_after=0):
    """runs: список абзацев; абзац = список (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        for text, size, color, bold in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Roboto"
    return tb


def _kpi_strip(slide, y, kpi: dict, fleet: dict):
    """Полоса KPI парка: плитки в один ряд."""
    tiles = [
        ("ТС в программе", _ru(fleet.get("vehicles")), f"с данными {_ru(fleet.get('with_data'))}"),
        ("Пробег", _ru(round((kpi.get("total_mileage_km") or 0) / 1000)) + " тыс км", "за период"),
        ("Топливо", _ru(round((kpi.get("total_fuel_l") or 0) / 1000)) + " тыс л", "фактический расход"),
        ("Моточасы", _ru(round(kpi.get("total_engine_hours") or 0)), "суммарно"),
        ("Холостой ход", _ru(round((kpi.get("fuel_idle_share") or 0) * 100)) + "%",
         "доля топлива простоя"),
        ("Пробег с превыш.", _ru(round((kpi.get("speeding_mileage_share") or 0) * 100, 1), 1) + "%",
         "от общего пробега"),
    ]
    n = len(tiles)
    gap = Inches(0.12)
    x0 = Inches(0.45)
    w = int((SLIDE_W - 2 * x0 - gap * (n - 1)) / n)
    h = Inches(0.92)
    for i, (label, value, note) in enumerate(tiles):
        x = x0 + i * (w + gap)
        _box(slide, x, y, w, h, fill=PANEL, line=LINE)
        _text(slide, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), h - Inches(0.16), [
            [(label.upper(), 8, MUTED, True)],
            [(value, 17, INK, True)],
            [(note, 8, FAINT, False)],
        ])


def _block(slide, x, y, w, h, title: str, lines: list, accent=BLUE):
    """Блок секции: заголовок + строки [(лейбл, значение, цвет-значения)]."""
    _box(slide, x, y, w, h, fill=WHITE, line=LINE)
    bar = _box(slide, x, y, Inches(0.055), h, fill=accent)
    bar.adjustments[0] = 0.0
    pad = Inches(0.16)
    _text(slide, x + pad, y + Inches(0.1), w - 2 * pad, Inches(0.3),
          [[(title.upper(), 10.5, INK, True)]])
    rows = []
    for label, value, color in lines:
        rows.append([(label + "  ", 9.5, MUTED, False), (str(value), 10.5, color or INK, True)])
    _text(slide, x + pad, y + Inches(0.42), w - 2 * pad, h - Inches(0.52), rows, space_after=3)


def _top_names(items, key, n=3, fmt=lambda it: ""):
    out = []
    for it in items[:n]:
        name = (it.get("name") or it.get("vehicle") or it.get(key) or "?")
        name = str(name)[:26]
        out.append((name, fmt(it), INK))
    return out


def build_slide(snap: dict, sections: list[str]) -> bytes:
    """Один слайд 16:9 из выбранных секций снапшота (снапшот уже отскоуплен)."""
    chosen = [s for s in SECTIONS if s in (sections or [])] or list(SECTIONS)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    _box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    org = (snap.get("orgs") or [{}])[0]
    org_name = org.get("name") or "Холдинг"
    period = (snap.get("period") or {}).get("label") or ""
    kpi = org.get("kpi") or {}
    fleet = snap.get("fleet") or {}

    # Шапка
    _text(slide, Inches(0.45), Inches(0.22), Inches(9.5), Inches(0.65), [
        [("Мониторинг автопарка · ", 20, INK, True), (org_name, 20, BLUE, True)],
        [(f"Период: {period}", 11, MUTED, False)],
    ])
    _text(slide, Inches(9.2), Inches(0.26), Inches(3.68), Inches(0.5), [
        [("Сформировано " + dt.datetime.now().strftime("%d.%m.%Y %H:%M"), 9, FAINT, False)],
        [("Источник: телеметрия Omnicomm", 9, FAINT, False)],
    ], align=PP_ALIGN.RIGHT)

    y = Inches(1.0)
    if "fleet" in chosen:
        _kpi_strip(slide, y, kpi, fleet)
        y = y + Inches(1.06)

    # Сетка блоков под остальные секции
    blocks = []
    if "economics" in chosen:
        eco = snap.get("economics") or {}
        buckets = sorted(eco.get("buckets") or [], key=lambda b: -(b.get("existing_kzt") or 0))
        lines = [("Выявленные потери", _mln(eco.get("total_existing_kzt")), RED),
                 ("Потенциал экономии", _mln(eco.get("total_potential_kzt")), GREEN)]
        for b in buckets[:3]:
            lines.append((str(b.get("label"))[:30], _mln(b.get("existing_kzt")), INK))
        blocks.append(("Деньги · где потери", lines, RED))
    if "speed" in chosen:
        recs = snap.get("recommendations") or []
        total_ep = sum(r.get("episodes") or 0 for r in recs)
        lines = [("ТС с превышениями", _ru(len(recs)), AMBER),
                 ("Эпизодов за период", _ru(total_ep), INK)]
        worst = sorted(recs, key=lambda r: -(r.get("episodes") or 0))
        lines += _top_names(worst, "terminal_id",
                            fmt=lambda r: f"{_ru(r.get('episodes'))} эп · до +{_ru(r.get('max_excess'))} км/ч")
        blocks.append(("Скоростной режим", lines, AMBER))
    if "fuel" in chosen:
        totals = (snap.get("fuel") or {}).get("totals") or {}
        lines = [("Заправки", _ru(round((totals.get("refuel_l") or 0) / 1000)) + " тыс л", GREEN),
                 ("Выдача (АТЗ)", _ru(round((totals.get("delivery_l") or 0) / 1000)) + " тыс л", INK),
                 ("Топливо простоя", _ru(round((kpi.get("fuel_idle_l") or 0) / 1000)) + " тыс л", AMBER),
                 ("Расход (подвижные)", _ru(kpi.get("mobile_fuel_per_100km"), 1) + " л/100 км", INK)]
        blocks.append(("Топливо", lines, GREEN))
    if "quality" in chosen:
        c = (snap.get("sensor_health") or {}).get("counts") or {}
        total = sum(c.values()) or 1
        lines = [("Онлайн (≤30 мин)", f"{_ru(c.get('online'))} · {round((c.get('online') or 0)*100/total)}%", GREEN),
                 ("Устарели (30 мин–24 ч)", _ru(c.get("stale")), AMBER),
                 ("Офлайн (>24 ч)", _ru(c.get("offline")), RED)]
        blocks.append(("Качество данных", lines, BLUE))
    if "maint" in chosen:
        mt = snap.get("maintenance") or {}
        c = mt.get("counts") or {}
        lines = [("Просрочено ТО", _ru(c.get("просрочено")), RED),
                 ("Ожидается", _ru(c.get("ожидается")), AMBER)]
        overdue = [i for i in (mt.get("items") or []) if i.get("status") == "просрочено"]
        lines += _top_names(overdue, "terminal_id",
                            fmt=lambda i: f"{_ru(i.get('km_since'))} км / {_ru(i.get('mh_since'))} мч")
        blocks.append(("Контроль ТО", lines, AMBER))
    if "tyres" in chosen:
        ty = snap.get("tyres") or {}
        c = ty.get("counts") or {}
        lines = [("Просрочена замена", _ru(c.get("просрочено")), RED),
                 ("Пора менять", _ru(c.get("пора менять")), AMBER),
                 ("Износ комплектов", _mln(ty.get("wear_kzt_total")), INK)]
        worst = [i for i in (ty.get("items") or []) if i.get("status") in ("просрочено", "пора менять")]
        lines += _top_names(worst, "terminal_id",
                            fmt=lambda i: f"{round((i.get('worn_share') or 0)*100)}% ресурса")
        blocks.append(("Шины (по пробегу)", lines, BLUE))

    if blocks:
        n = len(blocks)
        cols = n if n <= 3 else (2 if n == 4 else 3)
        rows_n = 1 if n <= 3 else 2
        gap = Inches(0.16)
        x0 = Inches(0.45)
        grid_w = SLIDE_W - 2 * x0
        bw = int((grid_w - gap * (cols - 1)) / cols)
        bottom = SLIDE_H - Inches(0.42)
        bh = int((bottom - y - gap * (rows_n - 1)) / rows_n)
        for i, (title, lines, accent) in enumerate(blocks):
            r, c = divmod(i, cols)
            _block(slide, x0 + c * (bw + gap), y + r * (Emu(bh) + gap), Emu(bw), Emu(bh),
                   title, lines, accent)

    # Футер
    _text(slide, Inches(0.45), SLIDE_H - Inches(0.34), Inches(12.4), Inches(0.26), [
        [("Данные: платформа мониторинга автопарка (Omnicomm) · аномалии — «требует проверки», "
          "суммы штрафов — только для дорог общего пользования (КоАП РК ст.592)", 8, FAINT, False)],
    ])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
