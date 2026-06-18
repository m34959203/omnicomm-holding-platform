#!/usr/bin/env python3
"""Рендер .pptx в PNG-превью без LibreOffice — через python-pptx + Pillow.

Назначение: визуальная проверка вёрстки слайдов (пустые зоны, наложения,
позиции таблиц/графиков) в средах, где нет LibreOffice Impress.

Рисует каждую фигуру по её реальному bbox: картинки вставляются как есть,
таблицы — сеткой с текстом, текст-боксы/автофигуры — заливкой и текстом.
Дополнительно печатает геометрический аудит: нижняя граница контента,
пустой «хвост» снизу и пересечения текстовых фигур.

Использование:
    PYTHONPATH=src python3 scripts/render_pptx_preview.py <file.pptx> [outdir]
"""
from __future__ import annotations

import sys
import os
import io

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

# 9525 EMU = 1 px при 96 dpi
EMU_PER_PX = 9525
FONT_REG = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

_FONT_CACHE: dict = {}


def _font(size_px: int, bold: bool):
    key = (size_px, bold)
    if key not in _FONT_CACHE:
        path = FONT_BOLD if bold else FONT_REG
        _FONT_CACHE[key] = ImageFont.truetype(path, max(6, size_px))
    return _FONT_CACHE[key]


def _px(emu) -> int:
    return int(emu / EMU_PER_PX) if emu is not None else 0


def _rgb(color, default=(30, 30, 30)):
    try:
        if color and color.type is not None and color.rgb is not None:
            c = color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return default


def _shape_fill_rgb(shape):
    try:
        fill = shape.fill
        if fill.type == 1:  # solid
            c = fill.fore_color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return None


def _wrap(draw, text, font, max_w):
    words = text.split()
    lines, cur = [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if draw.textlength(trial, font=font) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [""]


def _para_default(para):
    """Размер/жирность/цвет по первому run — фолбэк для пустых run."""
    r0 = para.runs[0] if para.runs else None
    size_pt, bold, color = 14, False, (30, 30, 30)
    if r0 is not None:
        if r0.font.size is not None:
            size_pt = r0.font.size.pt
        bold = bool(r0.font.bold)
        color = _rgb(r0.font.color)
    return size_pt, bold, color


def _draw_textframe(draw, tf, left, top, width, height):
    """Порановый инлайн-рендер с переносом — стиль каждого run сохраняется
    (важно для оценки реальной плотности: маркеры и тело разного цвета/веса)."""
    x0, y0 = left + 4, top + 3
    cursor_y = y0
    maxw = width - 8
    for para in tf.paragraphs:
        d_pt, d_bold, d_color = _para_default(para)
        align = para.alignment
        # собираем токены (слово, стиль) по всем run абзаца
        tokens = []
        runs = para.runs or []
        if not runs and para.text.strip():
            runs = None
        if runs is None:
            words = para.text.split()
            tokens = [(w, d_pt, d_bold, d_color) for w in words]
            line_pt = d_pt
        else:
            for r in runs:
                rpt = r.font.size.pt if r.font.size is not None else d_pt
                rbold = bool(r.font.bold)
                rcolor = _rgb(r.font.color, d_color)
                # сохраняем ведущие пробелы как «приклеенный» токен-разделитель
                parts = r.text.split(" ")
                for k, w in enumerate(parts):
                    if w == "" and k < len(parts) - 1:
                        continue
                    tokens.append((w, rpt, rbold, rcolor))
        if not tokens:
            cursor_y += int(d_pt * 96 / 72 * 0.5)
            continue
        # раскладка токенов в строки с переносом
        lines, cur, cur_w = [], [], 0
        for w, pt, b, c in tokens:
            fpx = int(pt * 96 / 72)
            f = _font(fpx, b)
            ww = draw.textlength(w + " ", font=f)
            if cur and cur_w + ww > maxw:
                lines.append(cur)
                cur, cur_w = [], 0
            cur.append((w, fpx, b, c, ww))
            cur_w += ww
        if cur:
            lines.append(cur)
        for ln in lines:
            line_h = max(fpx for _, fpx, _, _, _ in ln)
            line_w = sum(ww for *_, ww in ln)
            if align == PP_ALIGN.CENTER:
                x = left + (width - line_w) / 2
            elif align == PP_ALIGN.RIGHT:
                x = left + width - line_w - 4
            else:
                x = x0
            for w, fpx, b, c, ww in ln:
                draw.text((x, cursor_y), w, fill=c, font=_font(fpx, b))
                x += ww
            cursor_y += line_h + 4
    return cursor_y


# Поля ячейки PowerPoint по умолчанию ≈ 0.05" сверху/снизу.
CELL_VPAD = 10


def _cell_font(cell, header):
    p0 = cell.text_frame.paragraphs[0]
    r0 = p0.runs[0] if p0.runs else None
    size_px, bold, color = 13, header, (40, 40, 40)
    if r0 is not None:
        if r0.font.size is not None:
            size_px = int(r0.font.size.pt * 96 / 72)
        bold = bold or bool(r0.font.bold)
        color = _rgb(r0.font.color, color)
    return size_px, bold, color


def _table_row_heights(draw, tbl, col_w, set_h):
    """РЕАЛЬНАЯ высота строк, как их посчитает PowerPoint: max(заданная,
    контент с переносом). Длинные имена ТС переносятся → строка выше.
    Именно это раньше маскировалось (нормировкой в bbox) и прятало
    переполнение за футер."""
    heights = []
    rows = list(tbl.rows)
    for ri, row in enumerate(rows):
        need = 0
        for ci, cell in enumerate(row.cells):
            cw = col_w[ci] if ci < len(col_w) else 40
            txt = cell.text.strip()
            size_px, bold, _ = _cell_font(cell, ri == 0)
            font = _font(size_px, bold)
            wrap = cell.text_frame.word_wrap
            if txt and wrap is not False:
                n = len(_wrap(draw, txt, font, cw - 8))
            else:
                n = 1
            line_h = int(size_px * 1.3)
            need = max(need, n * line_h + 2 * CELL_VPAD)
        heights.append(max(_px(row.height), need))
    return heights


def _draw_table(draw, gframe, left, top, width):
    tbl = gframe.table
    cols = list(tbl.columns)
    col_w = [_px(c.width) for c in cols]
    tot = sum(col_w) or width
    col_w = [int(w * width / tot) for w in col_w]
    row_h = _table_row_heights(draw, tbl, col_w, _px(gframe.height))
    y = top
    for ri, row in enumerate(tbl.rows):
        x = left
        rh = row_h[ri]
        for ci, cell in enumerate(row.cells):
            cw = col_w[ci] if ci < len(col_w) else 40
            fill = _shape_fill_rgb(cell) or (
                (245, 247, 250) if ri == 0 else (255, 255, 255))
            draw.rectangle([x, y, x + cw, y + rh], fill=fill,
                           outline=(200, 205, 212), width=1)
            txt = cell.text.strip()
            if txt:
                size_px, bold, color = _cell_font(cell, ri == 0)
                font = _font(size_px, bold)
                wrap = cell.text_frame.word_wrap
                wlines = (_wrap(draw, txt, font, cw - 8)
                          if wrap is not False else [txt])
                ty = y + CELL_VPAD
                for wl in wlines:
                    draw.text((x + 4, ty), wl, fill=color, font=font)
                    ty += int(size_px * 1.3)
            x += cw
        y += rh
    return y  # реальный низ таблицы


def render_slide(slide, idx, w_px, h_px):
    img = Image.new("RGB", (w_px, h_px), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    content_bottom = 0
    text_boxes = []
    overflows = []  # (label, реальный_низ) — контент уходит за рабочую зону
    foot_top = int(h_px * 0.94)  # ~y6.4in: низ примечания/футера
    for shape in slide.shapes:
        left, top = _px(shape.left), _px(shape.top)
        width, height = _px(shape.width), _px(shape.height)
        if width <= 0 or height <= 0:
            continue
        st = shape.shape_type
        try:
            if st == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                pic = Image.open(io.BytesIO(blob)).convert("RGBA")
                pic = pic.resize((max(1, width), max(1, height)))
                img.paste(pic, (left, top), pic)
                content_bottom = max(content_bottom, top + height)
                text_boxes.append((left, top, width, height, "image"))
            elif st == MSO_SHAPE_TYPE.TABLE or shape.has_table:
                real_bottom = _draw_table(draw, shape, left, top, width)
                content_bottom = max(content_bottom, real_bottom)
                text_boxes.append((left, top, width, real_bottom - top, "table"))
                if real_bottom > foot_top + 6:
                    overflows.append(("таблица", real_bottom))
            else:
                fill = _shape_fill_rgb(shape)
                if fill is not None:
                    draw.rectangle([left, top, left + width, top + height],
                                   fill=fill)
                if shape.has_text_frame and shape.text_frame.text.strip():
                    real_bottom = _draw_textframe(draw, shape.text_frame,
                                                  left, top, width, height)
                    label = shape.text_frame.text[:30].replace("\n", " ")
                    ext = max(top + height, real_bottom)
                    content_bottom = max(content_bottom, ext)
                    text_boxes.append((left, top, width, ext - top, label))
                    # текст-блок переполняется, если реальный низ за зоной
                    # (футер — отдельный мелкий блок, его не считаем)
                    if real_bottom > foot_top + 6 and "·" not in label:
                        overflows.append((f"текст «{label}»", real_bottom))
                elif fill is not None:
                    content_bottom = max(content_bottom, top + height)
        except Exception as e:  # noqa: BLE001
            draw.rectangle([left, top, left + width, top + height],
                           outline=(220, 0, 0), width=1)
            draw.text((left + 2, top + 2), f"[{e}]", fill=(220, 0, 0),
                      font=_font(11, False))
    return img, content_bottom, text_boxes, overflows


def _overlaps(boxes):
    res = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            ax, ay, aw, ah, al = boxes[i]
            bx, by, bw, bh, bl = boxes[j]
            ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
            iy = max(0, min(ay + ah, by + bh) - max(ay, by))
            # реальная коллизия: заметное перекрытие по вертикали И горизонтали
            if iy > 14 and ix > 40 and ix * iy > 0.06 * min(aw * ah, bw * bh):
                res.append((al, bl, ix, iy))
    return res


def _max_empty_band(boxes, top_lim, bot_lim):
    """Самая высокая пустая горизонтальная полоса в теле слайда."""
    spans = sorted((max(top_lim, b[1]), min(bot_lim, b[1] + b[3]))
                   for b in boxes if b[1] + b[3] > top_lim and b[1] < bot_lim)
    if not spans:
        return bot_lim - top_lim, top_lim, bot_lim
    best, by0, by1 = 0, top_lim, top_lim
    cur = top_lim
    for s, e in spans:
        if s - cur > best:
            best, by0, by1 = s - cur, cur, s
        cur = max(cur, e)
    if bot_lim - cur > best:
        best, by0, by1 = bot_lim - cur, cur, bot_lim
    return best, by0, by1


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    path = sys.argv[1]
    outdir = sys.argv[2] if len(sys.argv) > 2 else "/tmp/pptx_preview"
    os.makedirs(outdir, exist_ok=True)
    prs = Presentation(path)
    w_px, h_px = _px(prs.slide_width), _px(prs.slide_height)
    imgs = []
    print(f"=== {os.path.basename(path)} — {len(prs.slides._sldIdLst)} "
          f"слайдов, {w_px}x{h_px}px ===")
    n_overflow = 0
    for idx, slide in enumerate(prs.slides, 1):
        img, cbottom, boxes, overflows = render_slide(slide, idx, w_px, h_px)
        p = os.path.join(outdir, f"slide_{idx:02d}.png")
        img.save(p)
        imgs.append(img)
        flags = []
        if overflows:
            n_overflow += 1
            flags.append(f"ПЕРЕПОЛНЕНИЕ за футер: {len(overflows)}")
        # пустая полоса внутри тела (между шапкой ~y110 и футером ~y675)
        band, by0, by1 = _max_empty_band(boxes, int(h_px * 0.16),
                                         int(h_px * 0.94))
        band_pct = 100 * band / h_px
        if band_pct > 20:
            flags.append(f"ПУСТАЯ ПОЛОСА {band_pct:.0f}% (y{by0}–{by1})")
        if cbottom < h_px * 0.25:
            flags.append("ПОЧТИ ПУСТОЙ слайд")
        ov = _overlaps(boxes)
        if ov:
            flags.append(f"НАЛОЖЕНИЙ: {len(ov)}")
        flag = ("  ⚠ " + "; ".join(flags)) if flags else "  ok"
        print(f"  s{idx:02d}: контент до y={cbottom}px,"
              f" макс.пустая полоса {band_pct:.0f}%{flag}")
        for lbl, by in overflows:
            print(f"        ↳ {lbl}: низ y={by}px > футер y={int(h_px*0.94)}px")
        for a, b, ix, iy in ov:
            print(f"        ↳ overlap «{a}» × «{b}» ({ix}x{iy}px)")
    # контактный лист
    cols = 4
    rows = (len(imgs) + cols - 1) // cols
    tw, th = w_px // 3, h_px // 3
    sheet = Image.new("RGB", (cols * tw + 20, rows * th + 20), (235, 235, 235))
    for i, im in enumerate(imgs):
        thumb = im.resize((tw - 6, th - 6))
        r, c = divmod(i, cols)
        sheet.paste(thumb, (10 + c * tw, 10 + r * th))
    sheet_path = os.path.join(outdir, "_contact_sheet.png")
    sheet.save(sheet_path)
    print(f"Превью: {outdir}/slide_*.png · контакт-лист: {sheet_path}")
    if n_overflow:
        print(f"ИТОГ: ПЕРЕПОЛНЕНИЕ на {n_overflow} слайд(ах) — контент за футером.")
    else:
        print("ИТОГ: переполнений за футер не найдено.")


if __name__ == "__main__":
    main()
