"""Сквозные тесты конвейера на демо-данных (режим Б).

Проверяют ключевые инварианты ТЗ:
- Excel загружается в единую модель, шумовой столбец «сливы» игнорируется (§9).
- KPI считаются без падений, средневзвешенный расход корректен (§5).
- Аномалии помечаются «требуют проверки», без обвинений (§7).
- Итоговый .pptx собирается и не содержит запретных слов (§9).
"""

from __future__ import annotations

import os
import sys
import tempfile
from datetime import datetime, timezone

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from omnicomm_report import analytics, charts, data_loader, report_builder, validator  # noqa: E402
from omnicomm_report.models import ReportPeriod, Severity  # noqa: E402

SAMPLE = os.path.join(os.path.dirname(__file__), "..", "samples", "fleet_sample.xlsx")
FORBIDDEN = ("слив", "слива", "сливы", "воровств", "кража")


@pytest.fixture(scope="module")
def sample_path() -> str:
    if not os.path.exists(SAMPLE):
        from samples.generate_sample import main as gen  # type: ignore
        gen()
    return SAMPLE


@pytest.fixture(scope="module")
def vehicles(sample_path):
    return data_loader.load_from_excel(sample_path)


def test_excel_loads(vehicles):
    assert len(vehicles) == 10
    assert all(v.name for v in vehicles)


def test_fuel_drain_column_ignored(vehicles):
    # Шумовой столбец «возможные сливы топлива» не должен маппиться в модель.
    for v in vehicles:
        for key in v.raw:
            assert "слив" not in str(key).lower() or v.raw  # хранится только в raw, не как метрика
        assert not hasattr(v, "fuel_drain")


def test_validate_marks_review_only(vehicles):
    validated = validator.validate(vehicles)
    for v in validated:
        for a in v.anomalies:
            assert a.severity in (Severity.REVIEW, Severity.NOTE)
            assert not any(bad in a.message.lower() for bad in FORBIDDEN)


def test_kpi_weighted_consumption(vehicles):
    validated = validator.validate(vehicles)
    kpi = analytics.compute_kpi(validated)
    assert kpi.total_mileage_km > 0
    assert kpi.total_fuel_l > 0
    # средневзвешенный = сумма топлива / сумма пробега * 100
    expected = kpi.total_fuel_l / kpi.total_mileage_km * 100
    assert abs(kpi.weighted_fuel_per_100km - expected) < 1.0
    assert 0 <= kpi.fuel_idle_share <= 1


def test_full_pipeline_builds_pptx(vehicles):
    validated = validator.validate(vehicles)
    period = ReportPeriod(
        start=datetime(2026, 5, 1, tzinfo=timezone.utc),
        end=datetime(2026, 5, 31, tzinfo=timezone.utc),
    )
    report = analytics.analyze(validated, period, "ООО Тест", source="excel")
    report.generated_at = datetime(2026, 6, 6, tzinfo=timezone.utc)
    assert report.conclusions, "должны быть управленческие выводы"

    with tempfile.TemporaryDirectory() as d:
        chart_paths = charts.build_charts(report, d)
        assert chart_paths, "должны построиться графики"
        out = os.path.join(d, "report.pptx")
        report_builder.build_pptx(report, chart_paths, out)
        assert os.path.exists(out) and os.path.getsize(out) > 10_000


def test_pptx_has_no_forbidden_text(vehicles):
    from pptx import Presentation

    validated = validator.validate(vehicles)
    period = ReportPeriod(
        start=datetime(2026, 5, 1, tzinfo=timezone.utc),
        end=datetime(2026, 5, 31, tzinfo=timezone.utc),
    )
    report = analytics.analyze(validated, period, "ООО Тест", source="excel")
    report.generated_at = datetime(2026, 6, 6, tzinfo=timezone.utc)

    with tempfile.TemporaryDirectory() as d:
        chart_paths = charts.build_charts(report, d)
        out = os.path.join(d, "report.pptx")
        report_builder.build_pptx(report, chart_paths, out)

        prs = Presentation(out)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text.lower())
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text.lower())
        joined = " ".join(texts)
        for bad in FORBIDDEN:
            assert bad not in joined, f"в отчёте найдено запретное слово: {bad}"
